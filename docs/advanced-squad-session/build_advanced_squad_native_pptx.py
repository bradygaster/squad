"""
Build Advanced Squad Patterns deck in the structure of
"Advanced Squad Patterns - Squad 0.10 - Tamir Dresher.pptx".

Per-slide chrome:
  - Top-right page number (small)
  - Eyebrow tag (ALL CAPS)
  - Title (bold)
  - One-sentence subtitle
  - 3-5 cards: icon + heading + 2-4 bullets   (or terminal/table for code/comparison)
  - Source citation footer (file:line)

4-pattern arc: Memory -> State -> Spawning -> Cross-squad.
Real captured outputs (memory CLI, two-layer demo, live upgrade) are pasted verbatim.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---- Palette ----
BG       = RGBColor(0x06, 0x11, 0x1F)
PANEL    = RGBColor(0x0F, 0x23, 0x38)
PANEL2   = RGBColor(0x0A, 0x18, 0x2A)
TEXT     = RGBColor(0xF8, 0xFA, 0xFC)
MUTED    = RGBColor(0xCB, 0xD5, 0xE1)
DIM      = RGBColor(0x94, 0xA3, 0xB8)
ACCENT   = RGBColor(0x38, 0xBD, 0xF8)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
YELLOW   = RGBColor(0xFB, 0xBF, 0x24)
RED      = RGBColor(0xFB, 0x71, 0x85)
VIOLET   = RGBColor(0xA7, 0x8B, 0xFA)
BORDER   = RGBColor(0x33, 0x41, 0x55)
HEADERBG = RGBColor(0x17, 0x25, 0x54)
LIVE     = RGBColor(0xF9, 0x73, 0x16)   # vivid orange \u2014 reserved for "switch to terminal now" slides

# ---- Geometry (inches) ----
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_X = 0.6
HEADER_TOP = 0.40
CONTENT_TOP = 2.05
CONTENT_W = SLIDE_W - 2 * MARGIN_X
CONTENT_H = SLIDE_H - CONTENT_TOP - 0.85

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
BLANK = prs.slide_layouts[6]

# -----------------------------------------------------------------
# Primitive helpers
# -----------------------------------------------------------------
def _set_text(tf, text, *, size=18, color=MUTED, bold=False, align=None, font="Segoe UI"):
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_rect(slide, x, y, w, h, *, fill=PANEL, line=BORDER, line_width=0.75, corner=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def add_text_box(slide, x, y, w, h, text, *, size=18, color=MUTED, bold=False,
                 align=None, anchor=None, font="Segoe UI"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top  = Inches(0.02); tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    _set_text(tf, text, size=size, color=color, bold=bold, align=align, font=font)
    return tb


def add_card(slide, x, y, w, h, icon, heading, bullets, *,
             accent=ACCENT, icon_size=22, head_size=17, body_size=12):
    """3-line card: icon row + heading + bulleted lines."""
    add_rect(slide, x, y, w, h, fill=PANEL, line=BORDER, corner=True)
    # Top accent strip
    add_rect(slide, x, y, w, 0.08, fill=accent, line=accent)
    # Icon
    add_text_box(slide, x + 0.18, y + 0.18, 0.55, 0.45, icon,
                 size=icon_size, color=accent, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # Heading
    add_text_box(slide, x + 0.75, y + 0.18, w - 0.85, 0.45, heading,
                 size=head_size, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # Bullet body
    body_y = y + 0.75
    body_h = h - 0.85
    tb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(body_y), Inches(w - 0.3), Inches(body_h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        bullet = p.add_run()
        bullet.text = "\u2022  "
        bullet.font.name = "Segoe UI"; bullet.font.size = Pt(body_size)
        bullet.font.color.rgb = accent; bullet.font.bold = True
        body = p.add_run()
        body.text = line
        body.font.name = "Segoe UI"; body.font.size = Pt(body_size)
        body.font.color.rgb = MUTED


def add_terminal(slide, x, y, w, h, lines, *, accent=ACCENT, label=None):
    add_rect(slide, x, y, w, h, fill=PANEL2, line=BORDER, corner=True)
    add_rect(slide, x, y, w, 0.05, fill=accent, line=accent)
    text_y = y + 0.20
    text_h = h - 0.25
    if label:
        add_text_box(slide, x + 0.18, y + 0.05, w - 0.3, 0.25, label,
                     size=10, color=accent, bold=True, font="Segoe UI Semibold")
        text_y = y + 0.32
        text_h = h - 0.38
    tb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(text_y), Inches(w - 0.3), Inches(text_h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = "Consolas"; run.font.size = Pt(11)
        if line.startswith("$") or line.startswith(">") or line.startswith("#"):
            run.font.color.rgb = TEXT
        elif line.startswith("//") or line.startswith("---"):
            run.font.color.rgb = DIM
        else:
            run.font.color.rgb = MUTED


def add_table(slide, x, y, w, h, headers, rows, *, accent=ACCENT, body_size=12):
    cols = len(headers)
    n = len(rows) + 1
    table = slide.shapes.add_table(n, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for c, hd in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = HEADERBG
        tf = cell.text_frame
        tf.margin_left = Inches(0.07); tf.margin_right = Inches(0.07)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        _set_text(tf, hd, size=13, color=TEXT, bold=True, font="Segoe UI Semibold")
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = PANEL
            tf = cell.text_frame
            tf.margin_left = Inches(0.07); tf.margin_right = Inches(0.07)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            _set_text(tf, val, size=body_size, color=MUTED)


def add_flow_strip(slide, x, y, w, h, nodes, *, accent=ACCENT):
    n = len(nodes)
    arrow_w = 0.30
    node_w = (w - arrow_w * (n - 1)) / n
    cur = x
    for i, label in enumerate(nodes):
        add_rect(slide, cur, y, node_w, h, fill=PANEL, line=accent, corner=True)
        add_text_box(slide, cur + 0.05, y, node_w - 0.1, h, label,
                     size=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI Semibold")
        cur += node_w
        if i < n - 1:
            add_text_box(slide, cur, y, arrow_w, h, "\u2192",
                         size=22, color=accent, bold=True, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)
            cur += arrow_w


def add_flow_strip_highlighted(slide, x, y, w, h, nodes, *, active=None, accent=ACCENT, dim=DIM):
    """Same as add_flow_strip but `active` (0-based index) is highlighted; others dimmed.
    Used for progressive-reveal workflow slides.
    """
    n = len(nodes)
    arrow_w = 0.30
    node_w = (w - arrow_w * (n - 1)) / n
    cur = x
    for i, label in enumerate(nodes):
        is_active = (i == active)
        line_color = accent if is_active else dim
        text_color = TEXT if is_active else dim
        fill_color = PANEL if is_active else PANEL2
        add_rect(slide, cur, y, node_w, h, fill=fill_color, line=line_color, corner=True)
        if is_active:
            add_rect(slide, cur, y, node_w, 0.06, fill=accent, line=accent)
        add_text_box(slide, cur + 0.05, y + 0.08, node_w - 0.1, h - 0.08, label,
                     size=11, color=text_color, bold=True, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI Semibold")
        cur += node_w
        if i < n - 1:
            arrow_color = accent if (i == active or i + 1 == active) else dim
            add_text_box(slide, cur, y, arrow_w, h, "\u2192",
                         size=22, color=arrow_color, bold=True, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)
            cur += arrow_w


# -----------------------------------------------------------------
# Slide chrome
# -----------------------------------------------------------------
PAGE_COUNT = {"n": 0}


def _background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False


def add_chrome(slide, *, eyebrow, title, subtitle="", source="", accent=ACCENT):
    """Standard content slide chrome matching the reference deck."""
    PAGE_COUNT["n"] += 1
    n = PAGE_COUNT["n"]
    _background(slide)
    # Header accent bar
    add_rect(slide, MARGIN_X, HEADER_TOP - 0.05, 0.40, 0.08, fill=accent, line=accent)
    # Page number (top right)
    add_text_box(slide, SLIDE_W - 1.2, HEADER_TOP - 0.10, 0.6, 0.4,
                 f"{n:02d}",
                 size=14, color=DIM, bold=True, align=PP_ALIGN.RIGHT, font="Consolas")
    # Eyebrow
    add_text_box(slide, MARGIN_X + 0.50, HEADER_TOP - 0.10, CONTENT_W - 1.5, 0.4,
                 eyebrow.upper(),
                 size=12, color=accent, bold=True, font="Segoe UI Semibold")
    # Title
    add_text_box(slide, MARGIN_X, HEADER_TOP + 0.32, CONTENT_W, 0.95,
                 title,
                 size=30, color=TEXT, bold=True)
    # Subtitle
    if subtitle:
        add_text_box(slide, MARGIN_X, HEADER_TOP + 1.25, CONTENT_W, 0.55,
                     subtitle,
                     size=15, color=MUTED)
    # Source footer
    if source:
        add_text_box(slide, MARGIN_X, SLIDE_H - 0.40, CONTENT_W, 0.30,
                     "Source: " + source,
                     size=10, color=DIM, font="Consolas")


def add_content_slide(*, eyebrow, title, subtitle="", source="", accent=ACCENT,
                      body, notes=""):
    slide = prs.slides.add_slide(BLANK)
    add_chrome(slide, eyebrow=eyebrow, title=title, subtitle=subtitle, source=source, accent=accent)
    body(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_section_divider(*, eyebrow, title, lede, accent=ACCENT, notes=""):
    PAGE_COUNT["n"] += 1
    n = PAGE_COUNT["n"]
    slide = prs.slides.add_slide(BLANK)
    _background(slide)
    # Big accent bar
    add_rect(slide, MARGIN_X, 2.6, 0.18, 2.5, fill=accent, line=accent)
    add_text_box(slide, SLIDE_W - 1.2, 0.30, 0.6, 0.4,
                 f"{n:02d}",
                 size=14, color=DIM, bold=True, align=PP_ALIGN.RIGHT, font="Consolas")
    add_text_box(slide, MARGIN_X + 0.45, 2.55, CONTENT_W, 0.5,
                 eyebrow.upper(),
                 size=14, color=accent, bold=True, font="Segoe UI Semibold")
    add_text_box(slide, MARGIN_X + 0.45, 3.05, CONTENT_W, 2.0,
                 title,
                 size=54, color=TEXT, bold=True)
    add_text_box(slide, MARGIN_X + 0.45, 4.95, CONTENT_W, 1.4,
                 lede,
                 size=20, color=MUTED)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


# -----------------------------------------------------------------
# Layout helper for "N cards across" rows
# -----------------------------------------------------------------
def cards_row(slide, y, h, items, *, gap=0.20):
    """Render N cards in a horizontal row. items = list of (icon, heading, bullets, accent_color)."""
    n = len(items)
    w = (CONTENT_W - gap * (n - 1)) / n
    for i, (icon, heading, bullets, color) in enumerate(items):
        x = MARGIN_X + i * (w + gap)
        add_card(slide, x, y, w, h, icon, heading, bullets, accent=color)


# =================================================================
# SLIDES
# =================================================================

# ---------- 01. Title ----------
def s01_title():
    slide = prs.slides.add_slide(BLANK)
    _background(slide)
    # Accent column on the left
    add_rect(slide, 0.55, 1.6, 0.20, 4.3, fill=ACCENT, line=ACCENT)
    # Eyebrow
    add_text_box(slide, 1.0, 1.55, 11.5, 0.5,
                 "ADVANCED SQUAD PATTERNS",
                 size=16, color=ACCENT, bold=True, font="Segoe UI Semibold")
    # Title
    add_text_box(slide, 1.0, 2.1, 11.5, 2.0,
                 "Memory \u2022 State \u2022 Spawning \u2022 Cross-squad \u2022 Monorepo",
                 size=44, color=TEXT, bold=True)
    # Lede
    add_text_box(slide, 1.0, 4.15, 11.5, 0.6,
                 "Once your AI Squad team is running, scale becomes an architecture problem.",
                 size=20, color=MUTED)
    # Start-here code card
    add_rect(slide, 1.0, 5.0, 7.5, 1.3, fill=PANEL2, line=ACCENT, corner=True)
    add_text_box(slide, 1.15, 5.05, 7.2, 0.35, "start here",
                 size=11, color=ACCENT, bold=True, font="Segoe UI Semibold")
    tb = slide.shapes.add_textbox(Inches(1.15), Inches(5.32), Inches(7.2), Inches(0.95))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(["npm i -g @bradygaster/squad-cli",
                              "squad --version   # 0.10.0",
                              "copilot --agent squad"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.name = "Consolas"; run.font.size = Pt(13); run.font.color.rgb = TEXT
    # Speaker block
    add_rect(slide, 8.8, 5.0, 3.8, 1.3, fill=PANEL, line=BORDER, corner=True)
    add_text_box(slide, 8.95, 5.05, 3.5, 0.5, "Tamir Dresher", size=18, color=TEXT, bold=True)
    add_text_box(slide, 8.95, 5.45, 3.5, 0.35, "\u2022 Contributor & maintainer",
                 size=11, color=MUTED)
    add_text_box(slide, 8.95, 5.70, 3.5, 0.35, "\u2022 Principal Software Engineer",
                 size=11, color=MUTED)
    add_text_box(slide, 8.95, 5.95, 3.5, 0.35, "\u2022 tamirdresher.com",
                 size=11, color=ACCENT)
    # Footer source
    add_text_box(slide, 0.55, 7.05, 12.3, 0.30,
                 "Source: package.json:3  \u2022  CHANGELOG.md:7",
                 size=10, color=DIM, font="Consolas")
    slide.notes_slide.notes_text_frame.text = (
        "Open with: this is the 'how do we scale?' session, not another hello world. "
        "Audience already ran squad init. Today is about what to do when you have many agents, "
        "long histories, multiple machines, and multiple squads."
    )


# ---------- 02. Session promise ----------
def s02_promise():
    def body(slide):
        cards_row(slide, CONTENT_TOP, 4.0, [
            ("\u2462", "How does it remember?",
             ["Durable team learnings",
              "Shared decisions",
              "Governed memory layers",
              "Context recovery"], ACCENT),
            ("\u2461", "How do we change team shape?",
             ["Presets",
              "SubSquads",
              "Mission-scoped child squads"], GREEN),
            ("\u2460", "How do squads talk?",
             ["Manifest discovery",
              "Delegated work as issues",
              "Distributed mesh"], YELLOW),
        ])
    add_content_slide(
        eyebrow="Session promise",
        title="The real questions start after init",
        subtitle="This is the 'how do we scale?' session \u2014 not another hello world.",
        body=body,
        accent=ACCENT,
        notes=(
            "Frame the three questions explicitly. The four patterns coming up are the answers."
        ),
    )


# ---------- 03. Roadmap ----------
def s03_roadmap():
    def body(slide):
        items = [
            ("1", "Pattern 1 \u2014 Memory model",
             ["What persists",
              "Who reads it",
              "Class \u2192 load guidance",
              "Retrieval uses both"], ACCENT),
            ("2", "Pattern 2 \u2014 State backends",
             ["local / orphan / two-layer",
              "Clean PRs by construction",
              "Live upgrade on this deck"], VIOLET),
            ("3", "Pattern 3 \u2014 Squad spawning",
             ["Presets bootstrap shape",
              "Worktrees isolate work",
              "HQ + child mission briefs"], GREEN),
            ("4", "Pattern 4 \u2014 Cross-squad",
             ["Manifest discovery",
              "Issue-based delegation",
              "Distributed mesh"], YELLOW),
            ("5", "Pattern 5 \u2014 Monorepo",
             ["Subfolder mode",
              "agentFileRoot vs teamRoot",
              "Workflows at git root"], VIOLET),
        ]
        n = len(items)
        gap = 0.18
        w = (CONTENT_W - gap * (n - 1)) / n
        h = 4.0
        for i, (icon, heading, bullets, color) in enumerate(items):
            x = MARGIN_X + i * (w + gap)
            add_card(slide, x, CONTENT_TOP, w, h, icon, heading, bullets,
                     accent=color, head_size=15, body_size=12)
    add_content_slide(
        eyebrow="Roadmap",
        title="Five patterns, one operating model",
        subtitle="Compress setup, isolate work, preserve context, share contracts \u2014 then a closing end-to-end trace.",
        body=body,
    )


# ---------- 04. Release snapshot ----------
def s04_release():
    def body(slide):
        items = [
            ("\u2192", "Cross-squad orchestration",
             ["squad discover", "squad delegate", "manifest schema"], ACCENT),
            ("\u2261", "State + memory hardening",
             ["state backends", "runtime state tools", "governed memory model"], VIOLET),
            ("\u25A4", "Preset system",
             ["SQUAD_HOME", "built-in presets", "init / list / show / apply / save"], GREEN),
            ("\u21BB", "Scale loops",
             ["squad loop", "fleet hybrid dispatch", "watch improvements"], YELLOW),
            ("\u2691", "Safety roster",
             ["Rai built-in", "fact-checker role", "memory diagnostics"], RED),
        ]
        n = len(items)
        w = (CONTENT_W - 0.18 * (n - 1)) / n
        h = 3.5
        for i, item in enumerate(items):
            icon, heading, bullets, color = item
            x = MARGIN_X + i * (w + 0.18)
            add_card(slide, x, CONTENT_TOP, w, h, icon, heading, bullets,
                     accent=color, head_size=14, body_size=12)
    add_content_slide(
        eyebrow="Release snapshot",
        title="Squad 0.10.0: what changed for this talk",
        subtitle="0.10.0 shipped June 7, 2026 \u2014 97 changesets across CLI and SDK (50 SDK + 71 CLI, deduplicated).",
        source="CHANGELOG.md:7-66",
        body=body,
    )


# ---------- 05. Vocabulary ----------
def s05_vocab():
    def body(slide):
        items = [
            ("\u2699", "Coordinator",
             ["Routes work", "Enforces handoffs", "Assembles results"], ACCENT),
            ("\u2638", "Agent",
             ["Specialist", "Charter + boundaries", "Owns history"], GREEN),
            ("\u270E", "Scribe",
             ["Silent logger", "Merges decisions", "Cross-agent updates"], ACCENT),
            ("\u21BB", "Ralph",
             ["Backlog loop", "Watch / heartbeat", "Keeps work moving"], YELLOW),
            ("\u2691", "Rai",
             ["RAI reviewer", "Traffic-light verdicts", "Guardrail, not wall"], RED),
            ("\u25A4", "State backend",
             ["Where mutable state lives", "Same agent protocol", "local / orphan / two-layer"], VIOLET),
        ]
        # 3x2 grid
        w = (CONTENT_W - 0.36) / 3
        h = 1.85
        for i, item in enumerate(items):
            icon, heading, bullets, color = item
            row, col = divmod(i, 3)
            x = MARGIN_X + col * (w + 0.18)
            y = CONTENT_TOP + row * (h + 0.18)
            add_card(slide, x, y, w, h, icon, heading, bullets,
                     accent=color, head_size=14, body_size=12)
    add_content_slide(
        eyebrow="Shared language",
        title="Vocabulary for the rest of the session",
        subtitle="If these words are clear, the advanced patterns become mechanical.",
        body=body,
    )


# ---------- 06. The scaling shift ----------
def s06_scaling_shift():
    def body(slide):
        # Five mechanisms; each names the user pain it solves and points to its later Pattern.
        items = [
            ("\u29C9", "Two-layer state",
             ["Clean PRs even when 5 agents",
              "are writing decisions",
              "",
              "\u2192 Pattern 2"], VIOLET),
            ("\u2691", "Worktrees",
             ["Run multiple missions",
              "without losing your branch",
              "",
              "\u2192 Pattern 3"], GREEN),
            ("\u2261", "SubSquads",
             ["Carve one repo into",
              "label-scoped lanes",
              "",
              "\u2192 Pattern 3"], GREEN),
            ("\u21C4", "Manifest + delegate",
             ["Hand work to another squad",
              "as a real GitHub issue",
              "",
              "\u2192 Pattern 4"], YELLOW),
            ("\u29DA", "Distributed mesh",
             ["Cross-machine peers",
              "via git pull / curl",
              "",
              "\u2192 Pattern 4"], YELLOW),
        ]
        cards_row(slide, CONTENT_TOP, 4.0, items)
        add_text_box(slide, MARGIN_X, CONTENT_TOP + 4.3, CONTENT_W, 0.5,
                     "Five separate mechanisms. Each solves a specific pain. Mix as needed \u2014 you don\u2019t need them all on day one.",
                     size=14, color=DIM, align=PP_ALIGN.CENTER, font="Segoe UI")
    add_content_slide(
        eyebrow="What scaling actually means",
        title="Scaling Squad isn\u2019t \u201Cmore agents\u201D \u2014 it\u2019s five concrete mechanisms",
        subtitle="Each one solves a real, separate pain. We\u2019ll cover them in the next four patterns.",
        body=body,
    )


# ====================== PART 1: MEMORY ======================

def s07_section_memory():
    add_section_divider(
        eyebrow="Pattern 1",
        title="\u2780 Memory that\ncompounds",
        lede="Agents stop asking the questions they already answered \u2014 if the right facts survive in the right layer.",
        accent=ACCENT,
        notes=("Section opener for the memory pattern. Anchor: 'agents getting fuller, not dumber.'"),
    )


def s08_three_layers():
    def body(slide):
        # Stacked horizontal bands. Top = most reusable / portable.
        # Bottom = most personal. Each band: icon+name (left), path+scope (middle), examples (right).
        layers = [
            {
                "icon":      "\u2698",
                "name":      "Skills",
                "subtitle":  "Reusable patterns \u2022 Team-wide \u2022 Portable",
                "path":      ".squad/skills/{name}/SKILL.md",
                "scope":     "All agents read on demand",
                "examples": ["\u201CHow to set up CI with GitHub Actions\u201D",
                             "\u201CDeploy to Azure Container Apps\u201D",
                             "\u201CReviewer protocol for security PRs\u201D"],
                "color":     YELLOW,
            },
            {
                "icon":      "\u270E",
                "name":      "Shared decisions",
                "subtitle":  "Team-wide rules \u2022 Every agent reads at spawn",
                "path":      ".squad/decisions.md",
                "scope":     "All agents \u2014 always loaded",
                "examples": ["\u201CUse PostgreSQL for primary storage\u201D",
                             "\u201CNo Friday deploys\u201D",
                             "\u201CAlways run SDK tests before merging\u201D"],
                "color":     ACCENT,
            },
            {
                "icon":      "\u2638",
                "name":      "Personal history",
                "subtitle":  "Per-agent domain memory \u2022 Owner-only",
                "path":      ".squad/agents/{name}/history.md",
                "scope":     "Only the owning agent reads",
                "examples": ["Kane: \u201CAuth uses JWT, refresh tokens in Redis\u201D",
                             "Dallas: \u201CUI tokens in src/theme/tokens.ts\u201D",
                             "Lambert: \u201CFlaky test pattern in checkout.spec\u201D"],
                "color":     GREEN,
            },
        ]
        # Layout — stacked bands across full content area
        band_h = 1.55
        gap    = 0.18
        for i, layer in enumerate(layers):
            y = CONTENT_TOP + i * (band_h + gap)
            color = layer["color"]
            # Band container
            add_rect(slide, MARGIN_X, y, CONTENT_W, band_h, fill=PANEL, line=BORDER, corner=True)
            # Left accent strip (the "depth" cue)
            add_rect(slide, MARGIN_X, y, 0.18, band_h, fill=color, line=color)
            # Layer label column (icon + name + subtitle)
            label_x = MARGIN_X + 0.40
            label_w = 3.4
            add_text_box(slide, label_x, y + 0.18, 0.65, 0.55, layer["icon"],
                         size=24, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, label_x + 0.65, y + 0.18, label_w - 0.7, 0.55, layer["name"],
                         size=20, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, label_x, y + 0.78, label_w, 0.35, layer["subtitle"],
                         size=11, color=MUTED, font="Segoe UI")
            # Path + scope column
            path_x = MARGIN_X + 0.40 + label_w
            path_w = 3.6
            add_text_box(slide, path_x, y + 0.22, path_w, 0.40, layer["path"],
                         size=12, color=color, bold=True, font="Consolas")
            add_text_box(slide, path_x, y + 0.65, path_w, 0.32, layer["scope"],
                         size=11, color=DIM, font="Segoe UI")
            # Examples column (right)
            ex_x = MARGIN_X + 0.40 + label_w + path_w + 0.20
            ex_w = CONTENT_W - (ex_x - MARGIN_X) - 0.15
            add_text_box(slide, ex_x, y + 0.10, ex_w, 0.30, "EXAMPLES",
                         size=9, color=color, bold=True, font="Segoe UI Semibold")
            tb = slide.shapes.add_textbox(Inches(ex_x), Inches(y + 0.40),
                                          Inches(ex_w), Inches(band_h - 0.50))
            tf = tb.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
            for j, ex in enumerate(layer["examples"]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_after = Pt(2)
                bullet = p.add_run(); bullet.text = "\u25B8  "
                bullet.font.name = "Segoe UI"; bullet.font.size = Pt(11)
                bullet.font.color.rgb = color; bullet.font.bold = True
                body_run = p.add_run(); body_run.text = ex
                body_run.font.name = "Segoe UI"; body_run.font.size = Pt(11)
                body_run.font.color.rgb = MUTED
    add_content_slide(
        eyebrow="Mental model",
        title="The three-layer memory model",
        subtitle="Not all memory belongs in the same place \u2014 most portable on top, most personal on the bottom.",
        source="docs/concepts/memory-and-knowledge.md:28-96, 174-200 (Skills section)",
        body=body,
    )


def s09_governed_classes():
    def body(slide):
        items = [
            ("TRANSIENT",       "Scratch / task-local",            "loads: NEVER",     RED),
            ("LOCAL",           "Safe local fact",                  "loads: ON-DEMAND", GREEN),
            ("DECISION",        "Candidate team truth",             "loads: ALWAYS",    ACCENT),
            ("POLICY",          "Future behavior rule",             "loads: ALWAYS",    YELLOW),
            ("COPILOT_MEMORY",  "External semantic",                "loads: ON-DEMAND", VIOLET),
            ("FORBIDDEN",       "Secrets / unsafe",                 "loads: NEVER",     RED),
        ]
        w = (CONTENT_W - 0.36) / 3
        h = 1.75
        for i, (name, what, guidance, color) in enumerate(items):
            row, col = divmod(i, 3)
            x = MARGIN_X + col * (w + 0.18)
            y = CONTENT_TOP + row * (h + 0.20)
            add_rect(slide, x, y, w, h, fill=PANEL, line=color, corner=True)
            add_rect(slide, x, y, w, 0.08, fill=color, line=color)
            add_text_box(slide, x + 0.2, y + 0.18, w - 0.4, 0.5,
                         name, size=15, color=color, bold=True, font="Consolas")
            add_text_box(slide, x + 0.2, y + 0.70, w - 0.4, 0.5,
                         what, size=13, color=TEXT)
            add_text_box(slide, x + 0.2, y + 1.20, w - 0.4, 0.4,
                         guidance, size=12, color=DIM, font="Consolas")
    add_content_slide(
        eyebrow="0.10 memory model",
        title="Memory has a class \u2014 and the class implies how it loads",
        subtitle="Decided at write time. Used at read time.",
        source="packages/squad-sdk/src/memory/index.ts:5-13, 428-440",
        body=body,
    )


def s09b_memory_tools():
    def body(slide):
        # 7-row catalog of governed memory tools.
        # Layout: each row is a band with: command (mono+bold), signature (mono+dim), purpose (text).
        tools = [
            ("squad memory classify", "(content) \u2192 { class, allowed, destination, loadGuidance, reason }",
             "Vet content. NO storage. Returns destination + load policy.", ACCENT),
            ("squad memory write",    "(content, --class?, --title?) \u2192 { stored, id, classification, path }",
             "classify \u2192 store \u2192 audit \u2192 replicate to registered providers.", GREEN),
            ("squad memory search",   "(query) \u2192 [{ id, class, loadGuidance, snippet, provider }, \u2026]",
             "Query active index. Query itself is classified; FORBIDDEN \u2192 [].", ACCENT),
            ("squad memory promote",  "(id, targetClass) \u2192 { stored, classification }",
             "Reclassify an existing entry. Cannot promote to FORBIDDEN or TRANSIENT.", VIOLET),
            ("squad memory delete",   "(id) \u2192 boolean",
             "Remove entry from index + emit audit event.", RED),
            ("squad memory audit",    "() \u2192 [event\u2026]",
             "Read the redacted forensics trail. Every tool emits exactly one.", YELLOW),
            ("squad memory provider", "() \u2192 { defaultProvider, realCopilotMemory: { available, reason } }",
             "Inspect what's actually configured. Fail-closed for missing providers.", DIM),
        ]
        # Geometry: 7 rows, each 0.55in tall
        row_h = 0.55
        gap   = 0.06
        # Header
        hdr_y = CONTENT_TOP
        add_rect(slide, MARGIN_X, hdr_y, CONTENT_W, 0.40, fill=HEADERBG, line=HEADERBG, corner=True)
        add_text_box(slide, MARGIN_X + 0.20, hdr_y + 0.04, 4.30, 0.32,
                     "COMMAND", size=11, color=ACCENT, bold=True, font="Segoe UI Semibold")
        add_text_box(slide, MARGIN_X + 4.55, hdr_y + 0.04, 4.30, 0.32,
                     "SIGNATURE", size=11, color=ACCENT, bold=True, font="Segoe UI Semibold")
        add_text_box(slide, MARGIN_X + 8.95, hdr_y + 0.04, CONTENT_W - 8.95 - 0.20, 0.32,
                     "PURPOSE", size=11, color=ACCENT, bold=True, font="Segoe UI Semibold")
        # Rows
        body_y0 = hdr_y + 0.40 + 0.08
        for i, (cmd, sig, purpose, color) in enumerate(tools):
            y = body_y0 + i * (row_h + gap)
            add_rect(slide, MARGIN_X, y, CONTENT_W, row_h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, MARGIN_X, y, 0.12, row_h, fill=color, line=color)
            add_text_box(slide, MARGIN_X + 0.25, y + 0.10, 4.25, row_h - 0.18, cmd,
                         size=12, color=color, bold=True, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, MARGIN_X + 4.55, y + 0.10, 4.30, row_h - 0.18, sig,
                         size=10, color=MUTED, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, MARGIN_X + 8.95, y + 0.10, CONTENT_W - 8.95 - 0.20, row_h - 0.18,
                         purpose, size=11, color=TEXT, anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI")
    add_content_slide(
        eyebrow="The tools",
        title="The governed memory tools \u2014 contract surface",
        subtitle="Seven CLI/SDK operations. Each one classifies, audits, and respects the class on the way in and out.",
        source="packages/squad-cli/src/cli/commands/memory.ts:177-285 \u2022 packages/squad-sdk/src/memory/index.ts (LocalMemoryStore)",
        body=body,
        notes=("Each `squad memory <op>` command maps 1:1 to a method on LocalMemoryStore. "
               "The next slides use these tools in real CLI output \u2014 audience now knows their signatures."),
    )


def s09c_classify_spotlight():
    def body(slide):
        # Honest reality of the classifier (verified by running it against the real SDK).
        # Key finding: the (CI|PR|build) regex is in BOTH FORBIDDEN_PATTERNS (line 368)
        # AND the TRANSIENT heuristic (line 610). FORBIDDEN scan wins, so the TRANSIENT
        # heuristic path is dead code for these inputs. TRANSIENT only fires via explicit
        # --class override. Verified 2026-06-09 by harnessing LocalMemoryStore.classify().
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 7.0, 4.4,
                     ["# 1. FORBIDDEN scan runs FIRST \u2014 beats every override",
                      "$ squad memory classify \"sk-1234567890abcdef...\"",
                      "  class: FORBIDDEN    allowed: FALSE      reason: access token",
                      "",
                      "# 2. FORBIDDEN ALSO catches CI/PR/build status (same regex appears",
                      "#    in both lists \u2014 FORBIDDEN scan wins, so TRANSIENT heuristic is",
                      "#    unreachable from heuristics today; only fires via explicit --class)",
                      "$ squad memory classify \"CI build failed at 14:32\"",
                      "  class: FORBIDDEN    allowed: FALSE      reason: transient CI/PR status",
                      "",
                      "# 3. POLICY: /^(always|never|must|do not)/i",
                      "$ squad memory classify \"Always deploy on Tuesdays\"",
                      "  class: POLICY       allowed: true       loadGuidance: ALWAYS",
                      "",
                      "# 4. DECISION: /decision|decided|adopt|standardize/",
                      "$ squad memory classify \"We decided to use PostgreSQL\"",
                      "  class: DECISION     allowed: true       loadGuidance: ALWAYS",
                      "",
                      "# 5. LOCAL fallthrough \u2014 no signal, no danger",
                      "$ squad memory classify \"hello\"",
                      "  class: LOCAL        allowed: true       loadGuidance: ON-DEMAND",
                      "",
                      "# 6. explicit override \u2014 caller can force a class",
                      "$ squad memory classify \"hello\" --class POLICY",
                      "  class: POLICY       allowed: true       loadGuidance: ALWAYS"],
                     accent=ACCENT, label="six paths \u2014 verified against the real SDK (run them yourself)")
        # Right column: the mental model card
        add_card(slide, MARGIN_X + 7.2, CONTENT_TOP, 4.9, 4.4, "\u29BF",
                 "What classify actually does",
                 ["Vets content. NO storage.",
                  "Returns the verdict the WRITE would get.",
                  "",
                  "Real evaluation order (memory/index.ts:594-621):",
                  "  1. FORBIDDEN scan \u2014 wins over everything",
                  "  2. \u26A0 Heuristics block (TRANSIENT first,",
                  "     then POLICY/DECISION/COPILOT_MEMORY,",
                  "     LOCAL fallthrough)",
                  "",
                  "Honest finding worth knowing:",
                  "the (CI|PR|build) regex sits in BOTH lists \u2014",
                  "FORBIDDEN catches it first, so the TRANSIENT",
                  "heuristic is unreachable for those inputs.",
                  "Belt-and-suspenders by design."],
                 accent=ACCENT, head_size=15, body_size=12)
    add_content_slide(
        eyebrow="Spotlight: classify",
        title="`squad memory classify` \u2014 the governance gate, with one honest surprise",
        subtitle="Vets content. No storage. The CI/PR/build pattern overlap means FORBIDDEN catches it before TRANSIENT ever runs.",
        source="packages/squad-sdk/src/memory/index.ts:360-372 (FORBIDDEN_PATTERNS) + :594-621 (classify)  \u2022  verified live 2026-06-09",
        body=body,
        notes=("Demo order matches the code: FORBIDDEN scan first (catches secrets AND CI/PR/build status), then heuristics in actual order. "
               "Anchor the audience on the honest finding: the same regex is in BOTH FORBIDDEN_PATTERNS and the TRANSIENT heuristic. "
               "FORBIDDEN wins because it runs first, so CI status content classifies as FORBIDDEN, not TRANSIENT. "
               "The TRANSIENT heuristic is effectively dead code for those inputs today \u2014 it only fires when callers pass --class TRANSIENT explicitly. "
               "This is the kind of architectural rough edge you call out honestly, not hide. (Verified by harnessing classify() against \"CI build failed at 14:32\" "
               "and seeing the FORBIDDEN result with reason='transient CI/PR status'.)"),
    )


def s10_retrieval():
    def body(slide):
        # Flow strip
        add_flow_strip(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 0.55,
                       ["query", "classify query", "reject FORBIDDEN",
                        "scan active index", "skip COPILOT_MEMORY*", "match",
                        "return + guidance"], accent=ACCENT)
        # Left: retrieval rules card
        add_card(slide, MARGIN_X, CONTENT_TOP + 0.90, 5.7, 3.3, "\u29BF",
                 "Retrieval rules (memory.search)",
                 ["Query is itself classified \u2014 FORBIDDEN \u2192 []",
                  "provider=copilot fails closed without a bridge",
                  "Only status='active' entries are scanned",
                  "COPILOT_MEMORY skipped unless host adapter present",
                  "Each result returns class + loadGuidance",
                  "Caller pins ALWAYS, surfaces ON-DEMAND, excludes ARCHIVE/NEVER"],
                 accent=ACCENT, head_size=14, body_size=12)
        # Right: verbatim search output
        add_terminal(slide, MARGIN_X + 5.9, CONTENT_TOP + 0.90, 6.2, 3.3,
                     ["$ squad memory search --query \"advanced demo snippet\"",
                      "search.complete  count=1  providers=local  elapsedMs=315",
                      "[",
                      "  {",
                      "    \"id\": \"e96f9d99-dbfc-492d-9e39-...\",",
                      "    \"class\": \"DECISION\",",
                      "    \"loadGuidance\": \"ALWAYS\",",
                      "    \"title\": \"Advanced demo snippet structure\",",
                      "    \"path\": \".squad/decisions/inbox/pao-...md\",",
                      "    \"snippet\": \"title: Advanced demo snippet...\",",
                      "    \"provider\": \"local\"",
                      "  }",
                      "]"],
                     label="real squad memory search output (run live)")
    add_content_slide(
        eyebrow="Retrieval",
        title="How classes shape what actually loads into context",
        subtitle="Class on write \u2192 load guidance on read. Memory becomes a queryable index.",
        source="packages/squad-sdk/src/memory/index.ts:863-966 \u2014 search() with class-aware filtering",
        body=body,
        notes=(
            "Walk the flow. Stress: search() classifies the QUERY first; FORBIDDEN search returns []. "
            "The result carries loadGuidance so the caller can pin (ALWAYS), surface (ON-DEMAND), or exclude (ARCHIVE/NEVER). "
            "Default class \u2192 guidance map: POLICY/DECISION=ALWAYS, LOCAL/COPILOT_MEMORY=ON-DEMAND, TRANSIENT/FORBIDDEN=NEVER."
        ),
    )


def s10a_write_architecture():
    def body(slide):
        # 5-layer stack diagram showing the full write path:
        # Agent \u2192 MCP tool \u2192 LocalMemoryStore (classify) \u2192 Provider(s) \u2192 Storage backend
        layers = [
            ("\u2638", "Agent",
             "\u201CI need to remember: Always run the SDK tests before merging.\u201D",
             "decides to persist a fact",
             ACCENT),
            ("\u29DA", "MCP tool surface",
             "squad_memory_write  /  squad_memory_classify  /  squad_memory_search",
             "exposed as MCP tools \u2014 agent calls by name with content + metadata",
             VIOLET),
            ("\u26A1", "LocalMemoryStore (orchestrator)",
             "classify(content) \u2192 { class, allowed, destination, loadGuidance }   then  storage.write()",
             "vets, classifies, renders frontmatter, emits audit \u2014 the gate every write goes through",
             GREEN),
            ("\u29C9", "MemoryProvider(s)  \u2014 plugin contract",
             "local  \u2022  hostInjectedCopilotAdapter  \u2022  MemPalace  \u2022  registered custom",
             "receive the safe, classified content. Only providers whose supportedClasses match get the write.",
             YELLOW),
            ("\u2698", "Storage backend (state-backend.ts)",
             "local FS  \u2022  orphan branch  \u2022  two-layer (branch + git note + audit)",
             "decides WHERE bytes land. Same agent protocol; backend swaps without code change.",
             RED),
        ]
        band_h = 0.92
        gap    = 0.10
        for i, (icon, name, line, sub, color) in enumerate(layers):
            y = CONTENT_TOP + i * (band_h + gap)
            add_rect(slide, MARGIN_X, y, CONTENT_W, band_h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, MARGIN_X, y, 0.18, band_h, fill=color, line=color)
            # Icon + name (left column)
            add_text_box(slide, MARGIN_X + 0.35, y + 0.10, 0.65, 0.42, icon,
                         size=22, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, MARGIN_X + 1.05, y + 0.10, 3.6, 0.42, name,
                         size=14, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            # Subtitle / role (left column, below name)
            add_text_box(slide, MARGIN_X + 0.35, y + 0.55, 4.3, 0.32, sub,
                         size=10, color=MUTED, font="Segoe UI")
            # Detail / signature (right column)
            add_text_box(slide, MARGIN_X + 4.85, y + 0.10, CONTENT_W - 5.00, band_h - 0.15, line,
                         size=11, color=TEXT, anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
            # Down-arrow between bands (except last)
            if i < len(layers) - 1:
                add_text_box(slide, MARGIN_X + (CONTENT_W - 0.40) / 2, y + band_h - 0.02, 0.40, gap + 0.04,
                             "\u2193", size=14, color=DIM, bold=True, align=PP_ALIGN.CENTER,
                             anchor=MSO_ANCHOR.MIDDLE)
    add_content_slide(
        eyebrow="The full write path",
        title="From \u201CI want to remember this\u201D to bytes on disk",
        subtitle="Five layers, one direction. Agent calls a tool; the store classifies; providers store; the backend decides where.",
        source="LocalMemoryStore.write() in memory/index.ts:667-861 \u2022 MemoryProvider interface at :134-143 \u2022 StateBackendStorageAdapter at state-backend.ts:791-840",
        body=body,
        accent=ACCENT,
        notes=("This slide answers the question \u2018what actually happens when an agent decides to remember something?\u2019. "
               "Top to bottom: agent invokes an MCP tool (squad_memory_write); the tool calls into the SDK\u2019s LocalMemoryStore; "
               "the store runs the SAME classify() gate covered earlier; if allowed, content + classification go to every "
               "provider whose supportedClasses match; the provider writes via the storage backend (which is where two-layer "
               "comes in). The next slide zooms the Provider layer."),
    )


def s10b_memory_providers():
    def body(slide):
        # LEFT \u2014 the contract (interface)
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.5, 4.2,
                     ["// packages/squad-sdk/src/memory/index.ts:134-143",
                      "interface MemoryProvider {",
                      "  readonly id:    string;       // \"local\", \"mempalace\", ...",
                      "  readonly name:  string;       // human label",
                      "  readonly supportedClasses:",
                      "    ReadonlyArray<MemoryClass>;",
                      "",
                      "  status():  Promise<MemoryProviderStatus>;",
                      "  write(req): Promise<...WriteResult>;",
                      "  search(q): Promise<...SearchResult[]>;",
                      "  delete(id): Promise<boolean>;",
                      "}",
                      "",
                      "// The classification flow runs ONCE in the store.",
                      "// Providers receive only safe, allowed content."],
                     accent=ACCENT, label="the storage plugin contract")
        # RIGHT \u2014 real providers shipped/expected
        providers = [
            ("\u25CF", "local",
             "Files in .squad/memory/{class}/   \u2022   ships by default",
             GREEN),
            ("\u25CB", "hostInjectedCopilotAdapter",
             "Opt-in slot for a host to bridge Copilot Memory",
             YELLOW),
            ("\u25CB", "real Copilot Memory (provider=copilot)",
             "Fails closed today \u2014 no callable API in the SDK",
             RED),
            ("\u25C6", "MemPalace (test double)",
             "Spatial loci, in-memory; for examples / tests",
             VIOLET),
            ("+ ", "Registered custom providers",
             "Any class implementing MemoryProvider; replicated to on write",
             DIM),
        ]
        # Compact provider rows
        x = MARGIN_X + 6.7
        w = CONTENT_W - 6.7
        row_h = 0.78
        for i, (glyph, name, desc, color) in enumerate(providers):
            y = CONTENT_TOP + i * (row_h + 0.05)
            add_rect(slide, x, y, w, row_h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, x, y, 0.12, row_h, fill=color, line=color)
            add_text_box(slide, x + 0.20, y + 0.06, 0.40, row_h - 0.12, glyph,
                         size=20, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, x + 0.65, y + 0.04, w - 0.80, 0.35, name,
                         size=12, color=TEXT, bold=True, font="Consolas")
            add_text_box(slide, x + 0.65, y + 0.38, w - 0.80, row_h - 0.42, desc,
                         size=11, color=MUTED, font="Segoe UI")
    add_content_slide(
        eyebrow="Storage plugin",
        title="Memory providers \u2014 what they are, and the ones that ship",
        subtitle="Same governed classification flows through every provider. The provider only decides WHERE bytes land.",
        source="packages/squad-sdk/src/memory/index.ts:134-143 \u2022 MemPalaceMemoryProvider at :156-227",
        body=body,
        notes=("Set up the audience BEFORE they see `squad memory provider` CLI output. "
               "Anchor: \u2018provider\u2019 = storage backend implementing a small interface (status/write/search/delete). "
               "Today local is the only one with real content; the hostInjectedCopilotAdapter slot exists so a host (VS Code, "
               "Codespaces) can bridge to Copilot Memory if they have an SDK; real provider=copilot fails closed because "
               "no callable Copilot Memory API is published yet. Squad refuses to fake it."),
    )


def s11_memory_tools_proof():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.5, 4.3,
                     ["$ squad memory provider --log-level info",
                      "provider.status.complete  defaultProvider=local",
                      "                          realCopilotConfigured=false",
                      "{",
                      "  \"defaultProvider\": \"local\",",
                      "  \"realCopilotMemory\": {",
                      "    \"available\": false,",
                      "    \"reason\": \"No callable Copilot API found.\"",
                      "  }",
                      "}",
                      "",
                      "$ squad memory write --class DECISION --approved \\",
                      "    --content \"Advanced demo snippets include ...\"",
                      "write.complete  stored=true  class=DECISION",
                      "                provider=local  loadGuidance=ALWAYS",
                      "                path=.squad/decisions/inbox/pao-...md"],
                     label="provider + write (verbatim)")
        add_terminal(slide, MARGIN_X + 6.7, CONTENT_TOP, 5.4, 4.3,
                     ["$ squad memory audit",
                      "audit.complete  count=2",
                      "[",
                      "  {",
                      "    \"timestamp\": \"2026-06-07T15:32:15Z\",",
                      "    \"action\": \"write\",",
                      "    \"id\": \"e96f9d99-...\",",
                      "    \"class\": \"DECISION\",",
                      "    \"actor\": \"pao\",",
                      "    \"provider\": \"local\"",
                      "  },",
                      "  {",
                      "    \"timestamp\": \"2026-06-07T15:32:18Z\",",
                      "    \"action\": \"search\",",
                      "    \"reason\": \"returned 1 result\"",
                      "  }",
                      "]"],
                     label="audit (verbatim)", accent=GREEN)
    add_content_slide(
        eyebrow="Verbatim CLI output",
        title="Memory tools are really called \u2014 and audited",
        subtitle="Local governed memory is proven. Remote Copilot Memory is not claimed.",
        source="packages/squad-cli/src/cli/commands/memory.ts:177-332 \u2014 provider, write, audit subcommands",
        body=body,
        accent=GREEN,
    )


def s12_session_evidence():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 2.6,
                     ["PR #1145 \u2014 paired real Copilot CLI A/B harness, isolated COPILOT_HOME:",
                      "",
                      "External Squad-using project demo:",
                      "  baseline                   :  0 memory diagnostic events",
                      "  memory-governance variant  : 10 memory diagnostic events",
                      "",
                      "Squad repo:",
                      "  baseline                   :  0 memory diagnostic events",
                      "  memory-governance variant  :  9 memory diagnostic events"],
                     accent=GREEN, label="paired A/B harness extract from PR #1145")
        add_text_box(slide, MARGIN_X, CONTENT_TOP + 2.85, CONTENT_W, 1.2,
                     "Same prompt, same model, same repo \u2014 only memory governance is switched on. "
                     "The signal jumps from zero to a measurable count of memory operations.",
                     size=15, color=MUTED)
    add_content_slide(
        eyebrow="Session evidence",
        title="Agents actually use the memory layer",
        subtitle="Paired A/B harness from PR #1145 \u2014 governance-on raises memory operations from 0 to N.",
        source="PR #1145 \u2014 paired real Copilot CLI A/B harness with isolated COPILOT_HOME",
        body=body,
        accent=GREEN,
    )


def s13_directive_capture():
    def body(slide):
        steps = [
            ("1", "You say",          "\u201cAlways run the SDK tests before merging.\u201d"),
            ("2", "Coordinator",      "Detects signals: always / never / must / do not"),
            ("3", "State write",      "Creates decisions/inbox/directive-*.md"),
            ("4", "Scribe merges",    "Deduplicates into decisions.md"),
            ("5", "Agents reload",    "Every future agent reads the shared rule"),
        ]
        n = len(steps)
        w = (CONTENT_W - 0.18 * (n - 1)) / n
        h = 3.6
        for i, (num, head, body_line) in enumerate(steps):
            x = MARGIN_X + i * (w + 0.18)
            add_rect(slide, x, CONTENT_TOP, w, h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, x, CONTENT_TOP, w, 0.08, fill=ACCENT, line=ACCENT)
            # Number badge
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                Inches(x + (w - 0.7) / 2), Inches(CONTENT_TOP + 0.30),
                Inches(0.7), Inches(0.7))
            badge.fill.solid(); badge.fill.fore_color.rgb = PANEL2
            badge.line.color.rgb = ACCENT; badge.line.width = Pt(1.5); badge.shadow.inherit = False
            tf = badge.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _set_text(tf, num, size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
            add_text_box(slide, x + 0.1, CONTENT_TOP + 1.15, w - 0.2, 0.5,
                         head, size=15, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
            add_text_box(slide, x + 0.1, CONTENT_TOP + 1.70, w - 0.2, 1.7,
                         body_line, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_content_slide(
        eyebrow="Step by step",
        title="Directive capture: the fastest way to teach the team",
        subtitle="Use natural language. The system turns it into durable team memory.",
        source="docs/.../concepts/memory-and-knowledge.md:103-150",
        body=body,
    )


def s14_memory_antipatterns():
    def body(slide):
        items = [
            ("\u26A0", "Secrets are never memory",
             ["Tokens, passwords, keys",
              "Raw customer data"], RED),
            ("\u23F1", "Transient data expires",
             ["CI blips, one-off logs",
              "Stale PR status",
              "Don't promote to DECISION"], YELLOW),
            ("\u2716", "Don't fake providers",
             ["Copilot Memory needs a real bridge",
              "provider=copilot fails closed",
              "Audit honest provider boundary"], RED),
            ("\u29C9", "Don't load everything",
             ["Archive old decisions",
              "Summarize long histories",
              "ALWAYS budget is small"], ACCENT),
        ]
        cards_row(slide, CONTENT_TOP, 4.0, items)
    add_content_slide(
        eyebrow="What NOT to persist",
        title="Memory anti-patterns",
        subtitle="The memory model is only useful if it stays safe, current, and intentional.",
        source="docs/.../features/memory.md:21-24, 93-109, 162-166",
        body=body,
    )


# ====================== PART 2: STATE ======================

def s15_section_state():
    add_section_divider(
        eyebrow="Pattern 2",
        title="\u2781 State that\nstays out of the PR",
        lede="Code and squad state are different workflows. Treat them like it.",
        accent=VIOLET,
        notes=("Section opener for state backends. Use Tamir Part 7/7b essay analogy."),
    )


def s16_state_vs_code():
    def body(slide):
        # Two halves: humans-review (LEFT) vs runtime-only (RIGHT)
        add_card(slide, MARGIN_X, CONTENT_TOP, 5.95, 3.5, "\u2696",
                 "Things humans review",
                 ["Code, tests, configs",
                  "Agent charters, team.md, routing.md",
                  "Shows in PR diffs",
                  "Shows in git log",
                  "Cadence: per feature \u2014 reviewed, gated"],
                 accent=ACCENT, head_size=16, body_size=13)
        add_card(slide, MARGIN_X + 6.18, CONTENT_TOP, 5.95, 3.5, "\u26A1",
                 "Things only the runtime cares about",
                 ["Every decision recorded",
                  "Every orchestration log line",
                  "Every agent's history.md",
                  "Changes dozens of times per session",
                  "On main \u2192 drowns git log + merge conflicts"],
                 accent=VIOLET, head_size=16, body_size=13)
        # Bottom thesis box
        add_rect(slide, MARGIN_X, CONTENT_TOP + 3.75, CONTENT_W, 1.05,
                 fill=PANEL2, line=VIOLET, corner=True)
        add_text_box(slide, MARGIN_X + 0.30, CONTENT_TOP + 3.85, CONTENT_W - 0.60, 0.40,
                     "Two-layer\u2019s job, in one sentence",
                     size=12, color=VIOLET, bold=True, font="Segoe UI Semibold")
        add_text_box(slide, MARGIN_X + 0.30, CONTENT_TOP + 4.20, CONTENT_W - 0.60, 0.55,
                     "Keep the things humans review on main (clean PRs) "
                     "and shove the things only the runtime cares about somewhere else (no noise).",
                     size=16, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_content_slide(
        eyebrow="The problem two-layer solves",
        title="Your team produces two completely different kinds of files",
        subtitle="One is for humans. One is for the runtime. They want different storage.",
        source="tamirdresher.com Part 7 / 7b \u2022 PR #1004 \u2022 state-backend.ts",
        body=body,
        accent=VIOLET,
    )


def s18_two_layer_thesis():
    def body(slide):
        # Three stacked horizontal layers \u2014 same visual pattern as the three-layer
        # memory slide. Storage | what it is | analogy.
        layers = [
            {
                "icon":    "\u29C8",
                "name":    "main",
                "what":    "Your normal branch",
                "analogy": "\u201CThe published book\u201D",
                "color":   ACCENT,
                "loc":     "refs/heads/main",
            },
            {
                "icon":    "\u270E",
                "name":    "squad-state",
                "what":    "A second branch with no shared history \u2014 you can't merge it. A parallel filesystem hiding in the same .git folder.",
                "analogy": "\u201CA separate author\u2019s journal \u2014 same library, never bound into the book\u201D",
                "color":   GREEN,
                "loc":     "refs/heads/squad-state (orphan)",
            },
            {
                "icon":    "\u29DA",
                "name":    "refs/notes/squad",
                "what":    "Tiny attachments stuck onto specific commits. They don't appear in git log by default.",
                "analogy": "\u201CSticky notes peeled onto specific pages of the book\u201D",
                "color":   YELLOW,
                "loc":     "refs/notes/squad",
            },
        ]
        band_h = 1.30
        gap    = 0.18
        for i, layer in enumerate(layers):
            y = CONTENT_TOP + i * (band_h + gap)
            color = layer["color"]
            add_rect(slide, MARGIN_X, y, CONTENT_W, band_h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, MARGIN_X, y, 0.18, band_h, fill=color, line=color)
            # Icon + name + ref location (left column)
            add_text_box(slide, MARGIN_X + 0.35, y + 0.20, 0.70, 0.55, layer["icon"],
                         size=28, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, MARGIN_X + 1.10, y + 0.10, 3.6, 0.45, layer["name"],
                         size=20, color=TEXT, bold=True, font="Consolas")
            add_text_box(slide, MARGIN_X + 1.10, y + 0.58, 3.6, 0.30, layer["loc"],
                         size=10, color=DIM, font="Consolas")
            # What it is (middle column)
            mid_x = MARGIN_X + 4.85
            mid_w = 4.4
            add_text_box(slide, mid_x, y + 0.10, mid_w, 0.30, "WHAT IT IS",
                         size=9, color=color, bold=True, font="Segoe UI Semibold")
            add_text_box(slide, mid_x, y + 0.40, mid_w, band_h - 0.50, layer["what"],
                         size=12, color=MUTED, font="Segoe UI")
            # Analogy (right column)
            ana_x = MARGIN_X + 9.40
            ana_w = CONTENT_W - (ana_x - MARGIN_X) - 0.15
            add_text_box(slide, ana_x, y + 0.10, ana_w, 0.30, "ANALOGY",
                         size=9, color=color, bold=True, font="Segoe UI Semibold")
            add_text_box(slide, ana_x, y + 0.40, ana_w, band_h - 0.50, layer["analogy"],
                         size=12, color=TEXT, font="Segoe UI")
    add_content_slide(
        eyebrow="Where \u201Csomewhere else\u201D actually is",
        title="Two-layer is three storages \u2014 same library, different bindings",
        subtitle="main is the book. squad-state is the journal. refs/notes/squad are sticky notes on specific pages.",
        source="PR #1004 \u2022 state-backend.ts: TwoLayerBackend (composes orphan + notes)",
        body=body,
        accent=VIOLET,
        notes=("Anchor the analogy hard: same library (same .git folder), different bindings. "
               "Orphan = no shared history with main, you can't merge it \u2014 it's a parallel filesystem. "
               "Notes = peeled onto specific pages, invisible in git log."),
    )


def s17_backends_compare():
    def body(slide):
        add_table(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 2.9,
                  ["", "Local", "Orphan branch", "Two-layer"],
                  [("Location",         ".squad/ files",        "squad-state branch",      "orphan branch + notes"),
                   ("PR cleanliness",   "Appears if committed", "Clean",                   "Clean"),
                   ("Branch-switch",    "Risk if uncommitted",  "Safe",                    "Safe"),
                   ("Team concurrency", "Filesystem / git merge", "Single-writer preferred", "Designed for teams")],
                  body_size=12)
        add_terminal(slide, MARGIN_X, CONTENT_TOP + 3.2, CONTENT_W, 1.45,
                     ["squad init --state-backend two-layer   # recommended for teams who want clean PRs + shared state",
                      "",
                      "# Preconditions:  inside a git repo  \u2022  `git remote -v` shows origin  \u2022  working tree clean",
                      "# Verify after:   squad doctor   (silent fallback to local if backend init fails \u2014 doctor surfaces it)"],
                     accent=VIOLET)
    add_content_slide(
        eyebrow="0.10 persistence model",
        title="State backends: where the memory lives",
        subtitle="Agents use the same state protocol; the backend decides where state is stored.",
        source="docs/.../features/state-backends.md:21-41, 63-156",
        body=body,
        accent=VIOLET,
    )


def s18_two_layer_thesis_old_REMOVED():
    """Replaced by the redesigned s18_two_layer_thesis above (kept stub to avoid
    accidentally re-wiring an old definition)."""
    pass


def s19_activation():
    def body(slide):
        add_flow_strip(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 0.6,
                       [".squad/config.json", "resolveStateBackend()",
                        "require git repo", "TwoLayerBackend", "orphan + notes"],
                       accent=VIOLET)
        add_terminal(slide, MARGIN_X, CONTENT_TOP + 0.95, 6.0, 2.4,
                     ["// .squad/config.json",
                      "{",
                      "  \"version\": 1,",
                      "  \"stateBackend\": \"two-layer\"",
                      "}"], accent=VIOLET, label="one config line")
        add_card(slide, MARGIN_X + 6.2, CONTENT_TOP + 0.95, 5.9, 2.4, "\u2699",
                 "Activation behavior",
                 ["Composes GitNotesBackend + OrphanBranchBackend",
                  "Reads come from orphan branch",
                  "Writes hit orphan + best-effort git note",
                  "Backend init failure \u2192 warn and fall back local"],
                 accent=VIOLET, head_size=15, body_size=12)
    add_content_slide(
        eyebrow="Step by step",
        title="How two-layer activates",
        subtitle="One config line. Real git plumbing under the hood.",
        source="packages/squad-sdk/src/state-backend.ts:1089-1186",
        body=body,
        accent=VIOLET,
    )


def s20_live_upgrade():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 4.0,
                     ["$ squad upgrade --state-backend two-layer",
                      "",
                      "Migrating state backend: local \u2192 two-layer",
                      "  \u2713 squad-state branch ready",
                      "  \u2713 migrated 18 state file(s) onto squad-state branch:",
                      "      .squad/decisions.md",
                      "      .squad/agents/{18 agents}/history.md",
                      "  \u2713 removed 18 stale working-tree file(s)",
                      "  \u2713 config.json updated: stateBackend = two-layer",
                      "",
                      "Installing squad sync hooks",
                      "  \u2713 pre-push, post-merge, post-rewrite, post-checkout, pre-commit, post-commit",
                      "",
                      "\u2713 Migration complete. Backend is now 'two-layer'."],
                     accent=LIVE, label="ran on this deck's worktree, today")
    add_content_slide(
        eyebrow="\u26A1 LIVE \u2014 RUN NOW",
        title="One command to switch this repo to two-layer",
        subtitle="Real migration, real hooks, real state branch \u2014 not a slide aspiration. (Run live during the talk.)",
        source="squad upgrade --state-backend two-layer  \u2022  packages/squad-cli/src/cli-entry.ts",
        body=body,
        accent=LIVE,
        notes=("\u26A1 SWITCH TO TERMINAL NOW \u2014 run `squad upgrade --state-backend two-layer` before advancing. "
               "Output should match the slide closely (file counts may differ). If it fails, fall back to narrating the slide."),
    )


def s21_live_state_proof():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.5, 4.2,
                     ["$ cat .squad/config.json",
                      "{",
                      "  \"version\": 1,",
                      "  \"stateBackend\": \"two-layer\"",
                      "}",
                      "",
                      "$ git ls-tree --name-only -r squad-state",
                      "README.md",
                      "agents/booster/history.md",
                      "agents/capcom/history.md",
                      "agents/control/history.md",
                      "agents/eecom/history.md",
                      "agents/fido/history.md",
                      "agents/flight/history.md",
                      "  ... 12 more agent histories",
                      "decisions.md"],
                     accent=GREEN, label="left this worktree clean")
        add_terminal(slide, MARGIN_X + 6.7, CONTENT_TOP, 5.4, 4.2,
                     ["$ git show squad-state:decisions.md",
                      "# Decisions",
                      "",
                      "> Team decisions that all agents",
                      "> must respect. Managed by Scribe.",
                      "",
                      "### Type safety \u2014 strict mode",
                      "**By:** CONTROL",
                      "**What:** `strict: true`, no",
                      "         `@ts-ignore` allowed.",
                      "**Why:** Types are contracts.",
                      "",
                      "### Hook-based governance",
                      "**By:** RETRO",
                      "**What:** Security/PII guards live",
                      "         in the hooks module.",
                      "**Why:** Prompts can be ignored;",
                      "         hooks are code."],
                     accent=ACCENT, label="real decisions, off the working tree")
    add_content_slide(
        eyebrow="\u26A1 LIVE \u2014 RUN NOW",
        title="The deck's own state is on the orphan branch",
        subtitle="Run these commands during the talk \u2014 they all return real content.",
        source="run `cat .squad/config.json` and `git ls-tree --name-only -r squad-state` live during the talk",
        body=body,
        accent=LIVE,
        notes=("\u26A1 SWITCH TO TERMINAL NOW \u2014 run the two commands shown on the left, then `git show squad-state:decisions.md` "
               "for the right panel. Three real commands; all return real content because two-layer is enabled on this repo."),
    )


def s21c_session_evidence():
    def body(slide):
        # Three storage rows showing what got written WHERE in a real session
        rows = [
            {
                "icon":   "\u29C8",
                "name":   "main",
                "count":  "1 commit",
                "detail": "4 Futurama charters + team.md / routing.md updates + .gitignore",
                "color":  ACCENT,
            },
            {
                "icon":   "\u270E",
                "name":   "squad-state",
                "count":  "7 commits",
                "detail": "3 directive files written to inbox  \u2192  Scribe merged them into decisions.md  \u2192  inbox files deleted  \u2192  orchestration logs for both Leela attempts  \u2192  Leela's history updated",
                "color":  GREEN,
            },
            {
                "icon":   "\u29DA",
                "name":   "refs/notes/squad",
                "count":  "1 snapshot update",
                "detail": "JSON blob containing decisions.md + all logs + Leela's history, pinned to commit e406cf7",
                "color":  YELLOW,
            },
        ]
        # Header strip
        hdr_y = CONTENT_TOP
        add_rect(slide, MARGIN_X, hdr_y, CONTENT_W, 0.40, fill=HEADERBG, line=HEADERBG, corner=True)
        add_text_box(slide, MARGIN_X + 0.20, hdr_y + 0.04, 3.5, 0.32,
                     "WHERE", size=11, color=ACCENT, bold=True, font="Segoe UI Semibold")
        add_text_box(slide, MARGIN_X + 4.30, hdr_y + 0.04, 2.5, 0.32,
                     "HOW MANY", size=11, color=ACCENT, bold=True, font="Segoe UI Semibold")
        add_text_box(slide, MARGIN_X + 7.10, hdr_y + 0.04, CONTENT_W - 7.10 - 0.20, 0.32,
                     "WHAT GOT WRITTEN", size=11, color=ACCENT, bold=True, font="Segoe UI Semibold")
        # Rows
        row_h = 1.20
        gap   = 0.14
        body_y0 = hdr_y + 0.40 + 0.10
        for i, row in enumerate(rows):
            y = body_y0 + i * (row_h + gap)
            color = row["color"]
            add_rect(slide, MARGIN_X, y, CONTENT_W, row_h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, MARGIN_X, y, 0.18, row_h, fill=color, line=color)
            # Where column \u2014 icon + name
            add_text_box(slide, MARGIN_X + 0.35, y + 0.20, 0.65, 0.45, row["icon"],
                         size=24, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, MARGIN_X + 1.05, y + 0.20, 3.20, 0.45, row["name"],
                         size=18, color=TEXT, bold=True, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
            # Count column
            add_text_box(slide, MARGIN_X + 4.30, y + 0.10, 2.7, row_h - 0.20, row["count"],
                         size=15, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            # Detail column
            detail_x = MARGIN_X + 7.10
            detail_w = CONTENT_W - 7.10 - 0.20
            add_text_box(slide, detail_x, y + 0.10, detail_w, row_h - 0.20, row["detail"],
                         size=11, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI")
        # Bottom takeaway
        cap_y = body_y0 + 3 * (row_h + gap) + 0.05
        add_rect(slide, MARGIN_X, cap_y, CONTENT_W, 0.50, fill=PANEL2, line=VIOLET, corner=True)
        add_text_box(slide, MARGIN_X + 0.30, cap_y + 0.08, CONTENT_W - 0.60, 0.35,
                     "main got 1 commit (the human-reviewable artifact).  "
                     "squad-state got 7 (the runtime\u2019s diary).  "
                     "refs/notes/squad got the frozen snapshot pinned to that main commit.",
                     size=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER, font="Segoe UI Semibold")
    add_content_slide(
        eyebrow="What happened in this session, in plain terms",
        title="The three storages, in real numbers",
        subtitle="Same audience question that built the deck \u2014 what does two-layer actually do in one session?",
        source="git log main \u2022 git log squad-state \u2022 git notes --ref=refs/notes/squad list",
        body=body,
        accent=VIOLET,
        notes=("Use this slide when the audience asks \u2018what does that look like for a real session?\u2019. "
               "main = 1 clean commit (the published book chapter). "
               "squad-state = 7 commits (the diary entries from the team's work). "
               "refs/notes/squad = the snapshot of what the team knew at that commit. "
               "If you check out e406cf7 with notes enabled, you get the full team context back without walking squad-state history."),
    )


# Three slides that zoom into a single real task: README on a feature branch.
# Source: files/two-layer-walkthrough.md \u2014 a captured trace from this session's repo.

def s21c2_branch_topology():
    def body(slide):
        # Two halves: LEFT = git history (real branches), RIGHT = parallel state world
        # Center divider with the key callout
        # LEFT panel
        add_rect(slide, MARGIN_X, CONTENT_TOP, 5.95, 4.4, fill=PANEL, line=ACCENT, corner=True)
        add_rect(slide, MARGIN_X, CONTENT_TOP, 5.95, 0.40, fill=ACCENT, line=ACCENT)
        add_text_box(slide, MARGIN_X + 0.20, CONTENT_TOP + 0.04, 5.5, 0.32,
                     "git history (real branches)",
                     size=12, color=TEXT, bold=True, font="Segoe UI Semibold")
        # ASCII diagram for left
        tb = slide.shapes.add_textbox(Inches(MARGIN_X + 0.25), Inches(CONTENT_TOP + 0.55),
                                      Inches(5.5), Inches(3.8))
        tf = tb.text_frame; tf.word_wrap = True
        for i, line in enumerate([
            "main:    \u25CF\u2500\u2500\u2500\u2500\u2500\u25CF\u2500\u2500\u2500\u2500\u2500\u25CF\u2500\u2500\u2500\u2500\u2500\u25CF",
            "         e406  24bf  361d  22f8  \u2190 HEAD",
            "         (4 commits \u2014 static identity)",
            "                                  \u2572",
            "                                   \u2572 checkout -b",
            "                                    \u2572",
            "                                     \u25CF  \u2190 feature/readme-\u2026",
            "                                     08da75c",
            "                                     README.md (28 lines)",
            "                                     \u2190 this is your PR diff",
        ]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run(); run.text = line
            run.font.name = "Consolas"; run.font.size = Pt(11)
            run.font.color.rgb = TEXT if "PR diff" in line or "HEAD" in line else MUTED
            if "feature/readme" in line or "08da75c" in line:
                run.font.color.rgb = ACCENT; run.font.bold = True
        # RIGHT panel
        add_rect(slide, MARGIN_X + 6.20, CONTENT_TOP, 5.95, 4.4, fill=PANEL, line=VIOLET, corner=True)
        add_rect(slide, MARGIN_X + 6.20, CONTENT_TOP, 5.95, 0.40, fill=VIOLET, line=VIOLET)
        add_text_box(slide, MARGIN_X + 6.40, CONTENT_TOP + 0.04, 5.5, 0.32,
                     "parallel state world (orphan + notes)",
                     size=12, color=TEXT, bold=True, font="Segoe UI Semibold")
        tb = slide.shapes.add_textbox(Inches(MARGIN_X + 6.45), Inches(CONTENT_TOP + 0.55),
                                      Inches(5.5), Inches(3.8))
        tf = tb.text_frame; tf.word_wrap = True
        for i, (line, color) in enumerate([
            ("squad-state:", DIM),
            ("\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u25CF\u2500\u25CF\u2500\u25CF\u2500\u25CF\u2500\u25CF\u2500\u25CF\u2500\u25CF  \u2190 squad-state HEAD", GREEN),
            ("(older sessions)  \u2570\u2500\u2500 7 commits from THIS task", MUTED),
            ("                     (1 leela-2 + 6 Scribe)", DIM),
            ("", DIM),
            ("refs/notes/squad:", DIM),
            ("\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u25CF\u2500\u25CF\u2500\u25CF\u2500\u25CF\u2500\u25CF\u2500\u25CF\u2500\u25CF  \u2190 refs/notes/squad", YELLOW),
            ("(older snapshots) \u2570\u2500\u2500 7 snapshot refreshes", MUTED),
            ("                     (one after each \u25CF, pinned", DIM),
            ("                      to main commit e406cf7)", DIM),
        ]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run(); run.text = line if line else " "
            run.font.name = "Consolas"; run.font.size = Pt(11)
            run.font.color.rgb = color
            if "HEAD" in line:
                run.font.bold = True
        # Center divider caption
        add_text_box(slide, MARGIN_X + 5.95, CONTENT_TOP + 1.95, 0.25, 0.5, "\u22EE",
                     size=20, color=DIM, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Bottom thesis
        cap_y = CONTENT_TOP + 4.55
        add_rect(slide, MARGIN_X, cap_y, CONTENT_W, 0.45, fill=PANEL2, line=BORDER, corner=True)
        add_text_box(slide, MARGIN_X + 0.30, cap_y + 0.08, CONTENT_W - 0.60, 0.30,
                     "No git relationship between these histories. main never sees squad-state. squad-state never sees main. Notes are the bridge.",
                     size=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER, font="Segoe UI Semibold")
    add_content_slide(
        eyebrow="The branch topology",
        title="Two parallel git universes \u2014 same repo, no shared history",
        subtitle="The PR lives on a feature branch. The team\u2019s memory lives on an orphan branch + notes ref. Both ride in the same .git folder.",
        source="files/two-layer-walkthrough.md \u2022 captured from squad-advanced-squad-session-slides repo",
        body=body,
        accent=VIOLET,
        notes=("Anchor: this is the deepest visual of the architecture in the deck. "
               "LEFT is what reviewers see. RIGHT is what the team remembers. "
               "They never merge \u2014 the orphan branch has no common ancestor with main. "
               "Notes are the only thing that bridges them, and only by attaching to specific main commits."),
    )


def s21c3_timeline():
    def body(slide):
        # 3-column timeline: TIME | LAYER | EVENT
        # Compressed to the most narrative-relevant rows.
        events = [
            ("15:32:49", "\u2328 user",        "\u201Clet\u2019s do README on a new branch\u201D", ACCENT),
            ("15:32:49", "\u25C6 coord",       "git checkout -b feature/readme-product-description", ACCENT),
            ("15:33:00", "\u25B6 spawn",       "task(name=leela, model=haiku) \u2192 leela-2", ACCENT),
            ("15:35:02", "\u2726 feature",     "\u25CF 08da75c   + README.md (28 lines)  \u2190 the PR diff", GREEN),
            ("15:35:18", "\u25A3 state",       "\u25CF 51b1966   + decisions/inbox/leela-readme-\u2026md  (22 lines: quote rationale)", VIOLET),
            ("15:35:19", "\u29DA note",        "\u25CF 4c546e7   refs/notes/squad refreshed", YELLOW),
            ("15:36:45", "\u25A3 state",       "\u25CF ed5d5c0   M decisions.md  (Scribe merged 2 inbox entries)", VIOLET),
            ("15:36:50", "\u25A3 state",       "\u25CF fe7cce0   D decisions/inbox/\u2026leela-readme\u2026  (inbox cleared)", VIOLET),
            ("15:37:06", "\u25A3 state",       "\u25CF 2f8f05a   + orchestration-log/\u2026-leela-2.md  (37 lines: forensic record)", VIOLET),
            ("15:37:08", "\u25A3 state",       "\u25CF ba586a4   + log/\u2026-readme-update.md  (session log)", VIOLET),
            ("15:37:10", "\u25A3 state",       "\u25CF 8661e3a   M agents/leela/history.md  (commit hash recorded)", VIOLET),
            ("15:37:11", "\u29DA note",        "\u25CF 2f21e96   \u2190 final snapshot of this work item", YELLOW),
        ]
        # Header
        hdr_y = CONTENT_TOP
        add_rect(slide, MARGIN_X, hdr_y, CONTENT_W, 0.36, fill=HEADERBG, line=HEADERBG, corner=True)
        add_text_box(slide, MARGIN_X + 0.20, hdr_y + 0.03, 1.30, 0.30, "TIME",
                     size=10, color=ACCENT, bold=True, font="Segoe UI Semibold")
        add_text_box(slide, MARGIN_X + 1.60, hdr_y + 0.03, 2.0, 0.30, "LAYER",
                     size=10, color=ACCENT, bold=True, font="Segoe UI Semibold")
        add_text_box(slide, MARGIN_X + 3.70, hdr_y + 0.03, CONTENT_W - 3.90, 0.30, "EVENT",
                     size=10, color=ACCENT, bold=True, font="Segoe UI Semibold")
        # Rows
        row_h = 0.34
        gap   = 0.02
        body_y0 = hdr_y + 0.36 + 0.04
        for i, (time, layer, event, color) in enumerate(events):
            y = body_y0 + i * (row_h + gap)
            add_rect(slide, MARGIN_X, y, CONTENT_W, row_h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, MARGIN_X, y, 0.10, row_h, fill=color, line=color)
            add_text_box(slide, MARGIN_X + 0.20, y + 0.03, 1.30, row_h - 0.06, time,
                         size=10, color=DIM, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, MARGIN_X + 1.60, y + 0.03, 2.0, row_h - 0.06, layer,
                         size=10, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI Semibold")
            add_text_box(slide, MARGIN_X + 3.70, y + 0.03, CONTENT_W - 3.90, row_h - 0.06, event,
                         size=10, color=TEXT, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
    add_content_slide(
        eyebrow="Timeline of a single task",
        title="One README change \u2014 1 feature commit, 7 state commits, 7 note refreshes",
        subtitle="2 minutes 22 seconds, real timestamps, real commit SHAs. Same repo. Three layers updating in lockstep.",
        source="files/two-layer-walkthrough.md \u2022 git log feature/readme-\u2026 + git log squad-state",
        body=body,
        accent=VIOLET,
        notes=("Walk the timeline top to bottom. Notice the rhythm: every squad-state commit (\u25A3) is followed within seconds "
               "by a notes refresh (\u29DA). The feature branch (\u2726) gets exactly ONE commit. Each color in the LAYER "
               "column maps to the topology slide before this one. Total task time: 2m22s; PR contribution: 1 commit / 28 lines."),
    )


def s21c4_landings():
    def body(slide):
        # 3 columns: PR (cyan), squad-state (violet), notes (yellow)
        cols = [
            {
                "icon": "\u2726",
                "title": "feature branch",
                "count": "1 commit  /  1 file  /  28 lines",
                "color": GREEN,
                "rows": [
                    ("08da75c", "+ README.md  (28 lines)"),
                    ("",        "title \u2022 quote \u2022 description"),
                    ("",        "platforms \u2022 status \u2022 features"),
                    ("",        "footer"),
                    ("",        ""),
                    ("",        "\u2190 the only thing a code reviewer sees"),
                ],
            },
            {
                "icon": "\u25A3",
                "title": "squad-state (orphan)",
                "count": "7 commits, invisible in the PR",
                "color": VIOLET,
                "rows": [
                    ("51b1966", "+ inbox/leela-readme-\u2026md  (rationale)"),
                    ("ed5d5c0", "M decisions.md  (Scribe merged 2)"),
                    ("58b61bd", "D inbox/\u2026copilot-directive\u2026"),
                    ("fe7cce0", "D inbox/\u2026leela-readme\u2026"),
                    ("2f8f05a", "+ orchestration-log/\u2026leela-2.md"),
                    ("ba586a4", "+ log/\u2026readme-update.md"),
                    ("8661e3a", "M agents/leela/history.md"),
                ],
            },
            {
                "icon": "\u29DA",
                "title": "refs/notes/squad",
                "count": "7 refreshes, one note object",
                "color": YELLOW,
                "rows": [
                    ("2f21e96", "pinned to main commit e406cf7"),
                    ("",        ""),
                    ("",        "JSON snapshot containing:"),
                    ("",        "  \u2022 decisions.md  (5 active)"),
                    ("",        "  \u2022 orchestration-log/\u2026leela-2.md"),
                    ("",        "  \u2022 log/\u2026readme-update.md"),
                    ("",        "  \u2022 agents/leela/history.md"),
                ],
            },
        ]
        n = len(cols)
        gap = 0.20
        w = (CONTENT_W - gap * (n - 1)) / n
        h = 4.0
        for i, col in enumerate(cols):
            x = MARGIN_X + i * (w + gap)
            color = col["color"]
            add_rect(slide, x, CONTENT_TOP, w, h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, x, CONTENT_TOP, w, 0.08, fill=color, line=color)
            # Icon + title
            add_text_box(slide, x + 0.15, CONTENT_TOP + 0.15, 0.55, 0.45, col["icon"],
                         size=22, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, x + 0.75, CONTENT_TOP + 0.15, w - 0.85, 0.45, col["title"],
                         size=15, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
            # Count line
            add_text_box(slide, x + 0.15, CONTENT_TOP + 0.65, w - 0.30, 0.30, col["count"],
                         size=11, color=color, bold=True, font="Segoe UI Semibold")
            # Rows (commit SHA + content)
            for j, (sha, content) in enumerate(col["rows"]):
                row_y = CONTENT_TOP + 1.05 + j * 0.40
                if sha:
                    add_text_box(slide, x + 0.15, row_y, 0.95, 0.32, sha,
                                 size=10, color=color, bold=True, font="Consolas")
                add_text_box(slide, x + 1.10 if sha else x + 0.15, row_y, w - (1.25 if sha else 0.30), 0.32, content,
                             size=10, color=MUTED if sha else TEXT, font="Consolas")
        # Bottom takeaway
        cap_y = CONTENT_TOP + 4.15
        add_rect(slide, MARGIN_X, cap_y, CONTENT_W, 0.55, fill=PANEL2, line=VIOLET, corner=True)
        add_text_box(slide, MARGIN_X + 0.30, cap_y + 0.10, CONTENT_W - 0.60, 0.40,
                     "Clone fresh \u2192 checkout main \u2192 `git notes --ref=squad show e406cf7`  \u2192  full team context reconstructed without walking 7 squad-state commits.",
                     size=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    add_content_slide(
        eyebrow="What landed where",
        title="One human-readable PR \u2014 14 invisible commits backing it up",
        subtitle="The PR reviewer sees 1 file, 28 lines. The team\u2019s memory has 7 state commits + 7 note refreshes \u2014 recoverable forever, never in the PR.",
        source="files/two-layer-walkthrough.md  \u2022  git show 08da75c \u2022 git log squad-state \u2022 git notes --ref=squad list",
        body=body,
        accent=VIOLET,
        notes=("This is the slide where the architecture pays off. PR is clean; team didn\u2019t lose any context. "
               "The bottom callout is the punchline: any future session/agent can recover the full team brain by pinning to one main commit."),
    )


def s21b_steady_state_safety_net():
    def body(slide):
        # Six hooks, each with a one-line WHY
        hooks = [
            ("pre-push",      "auto-push squad-state branch + refs/notes/squad*",
             "Teammates get your decisions when you push code \u2014 no separate sync step.",
             ACCENT),
            ("post-merge",    "fast-forward local squad-state from remote after pull",
             "You see the team\u2019s latest decisions immediately after `git pull`.",
             ACCENT),
            ("post-rewrite",  "same, after rebase",
             "Rebase doesn\u2019t leave squad-state drifting behind remote.",
             ACCENT),
            ("post-checkout", "fetch squad-state on branch switch",
             "State follows code context \u2014 switching branches refreshes the team brain.",
             ACCENT),
            ("pre-commit",    "REFUSE commits that stage .squad mutable state",
             "Defense-in-depth: prevents accidental leak of state into the working tree.",
             RED),
            ("post-commit",   "run `squad sync --quiet`",
             "Pending state pushed onto squad-state branch \u2014 your decisions don\u2019t sit only on your machine.",
             GREEN),
        ]
        # 6 rows, compact layout
        row_h = 0.70
        gap   = 0.08
        # Header
        hdr_y = CONTENT_TOP
        add_rect(slide, MARGIN_X, hdr_y, CONTENT_W, 0.40, fill=HEADERBG, line=HEADERBG, corner=True)
        add_text_box(slide, MARGIN_X + 0.20, hdr_y + 0.04, 2.6, 0.32,
                     "HOOK", size=12, color=ACCENT, bold=True, font="Segoe UI Semibold")
        add_text_box(slide, MARGIN_X + 2.95, hdr_y + 0.04, 4.5, 0.32,
                     "WHAT IT DOES", size=12, color=ACCENT, bold=True, font="Segoe UI Semibold")
        add_text_box(slide, MARGIN_X + 7.60, hdr_y + 0.04, CONTENT_W - 7.80, 0.32,
                     "WHY IT EXISTS", size=12, color=ACCENT, bold=True, font="Segoe UI Semibold")
        # Rows
        body_y0 = hdr_y + 0.40 + 0.08
        for i, (name, what, why, color) in enumerate(hooks):
            y = body_y0 + i * (row_h + gap)
            add_rect(slide, MARGIN_X, y, CONTENT_W, row_h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, MARGIN_X, y, 0.12, row_h, fill=color, line=color)
            add_text_box(slide, MARGIN_X + 0.22, y + 0.10, 2.65, row_h - 0.20, name,
                         size=13, color=color, bold=True, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, MARGIN_X + 2.95, y + 0.10, 4.5, row_h - 0.20, what,
                         size=12, color=MUTED, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, MARGIN_X + 7.60, y + 0.10, CONTENT_W - 7.80, row_h - 0.20, why,
                         size=12, color=TEXT, anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI")
    add_content_slide(
        eyebrow="The plumbing",
        title="Two-layer installs 6 git hooks \u2014 here\u2019s why each one",
        subtitle="`squad upgrade --state-backend two-layer` wires these. Each hook closes a specific failure mode.",
        source="install-hooks.ts:31-142, 286 \u2022 docs PR #1227 (docs only list 4 today)",
        body=body,
        accent=VIOLET,
        notes=("Walk through the rationale top to bottom. The first four are about KEEPING TEAMMATES IN SYNC \u2014 squad-state "
               "is useless if it only lives on your machine. The last two (pre-commit + post-commit) are the steady-state safety net: "
               "pre-commit STOPS leaks, post-commit AUTO-SYNCS. Without all six, two-layer would be theoretically clean but operationally noisy."),
    )


def s21d_pre_commit_fires():
    def body(slide):
        # Left \u2014 the failure scenario
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 7.0, 4.2,
                     ["$ git commit -am \"refactor: rename auth module\"",
                      "",
                      "\u26A0 squad pre-commit: refusing to commit",
                      "  two-layer state into the working tree.",
                      "  These paths belong on the 'squad-state'",
                      "  orphan branch, not in your normal commits:",
                      "",
                      "    .squad/decisions.md",
                      "    .squad/agents/scribe/history.md",
                      "    .squad/agents/leela/history.md",
                      "",
                      "  Use 'git restore --staged <path>' to unstage,",
                      "  or set SQUAD_SYNC_ACTIVE=1 to bypass.",
                      "",
                      "$ echo $LASTEXITCODE",
                      "1                              # commit refused"],
                     accent=RED, label="what the user sees")
        # Right \u2014 why it happened
        add_card(slide, MARGIN_X + 7.2, CONTENT_TOP, 4.9, 4.2, "\u26A0",
                 "Why those files reappeared",
                 ["The migration deleted them from disk \u2014",
                  "but SOMETHING wrote them back.",
                  "",
                  "Usual suspects:",
                  "  \u2022 an agent ran and wrote directly via fs",
                  "    (bypassing StateBackend)",
                  "  \u2022 an editor save touched the file",
                  "  \u2022 a legacy code path that still uses",
                  "    the local FS writer",
                  "",
                  "The hook is the backstop \u2014 it catches the leak",
                  "before it lands on main."],
                 accent=RED, head_size=14, body_size=12)
    add_content_slide(
        eyebrow="When the safety net fires",
        title="The pre-commit hook in action \u2014 and why it fires",
        subtitle="You\u2019ll see this. It\u2019s expected. The next slide is what you do about it.",
        source="install-hooks.ts:110-127 \u2022 PR #1229 prevents this with conditional .gitignore",
        body=body,
        accent=RED,
        notes=("Real audience moment. The user who triggered this whole series saw it after `git commit -am`. "
               "Hammer the WHY: state files reappear because not every write path goes through StateBackend yet. "
               "The hook is intentional defense-in-depth, not a bug."),
    )


def s21e_safety_net_recovery():
    def body(slide):
        # Left \u2014 recovery flow (the right thing to do)
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 7.0, 4.2,
                     ["# 1. unstage the state files",
                      "$ git restore --staged \\",
                      "    .squad/decisions.md \\",
                      "    .squad/agents/*/history.md",
                      "",
                      "# 2. verify the orphan branch already has them",
                      "$ git show squad-state:decisions.md",
                      "  (you should see your team\u2019s latest decisions)",
                      "",
                      "# 3. remove the stale working-tree copies",
                      "$ Remove-Item .squad/decisions.md, \\",
                      "              .squad/agents/*/history.md",
                      "",
                      "# 4. commit \u2014 hook is satisfied",
                      "$ git commit -m \"refactor: rename auth module\"",
                      "  \u2713 post-commit hook syncs squad-state to remote"],
                     accent=GREEN, label="the recovery flow \u2014 use this 99% of the time")
        # Right \u2014 the escape hatch
        add_card(slide, MARGIN_X + 7.2, CONTENT_TOP, 4.9, 4.2, "\u26A1",
                 "Escape hatch (don\u2019t use routinely)",
                 ["INLINE ONLY:",
                  "  SQUAD_SYNC_ACTIVE=1 git commit -m '...'",
                  "",
                  "\u26A0 NEVER `export` in .bashrc / .zshrc / profile",
                  "   \u2014 that PERMANENTLY disables the guard for",
                  "     every repo on your machine. Hard to debug.",
                  "",
                  "Use `git restore --staged` instead 99% of the time.",
                  "",
                  "What this flag does:",
                  "  \u2022 Bypasses the pre-commit guard",
                  "  \u2022 Commits state INTO the working tree \u2192",
                  "    defeats the clean-PR promise of two-layer",
                  "  \u2022 Internal Squad commands set it themselves",
                  "    (hence the recursion guard in every hook)"],
                 accent=YELLOW, head_size=14, body_size=11)
    add_content_slide(
        eyebrow="Recovery",
        title="Recovery flow \u2014 four commands, you\u2019re done",
        subtitle="Almost always: unstage, verify, delete, commit. The escape hatch exists but isn\u2019t the answer.",
        source="install-hooks.ts:110-142 \u2022 see PR #1229 which makes this rarely needed",
        body=body,
        accent=GREEN,
    )


def s22_two_layer_demo_promotion():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.6, 4.2,
                     ["$ git notes --ref=squad-flight add -m '...promote_to_permanent: true'",
                      "$ git notes --ref=squad-flight add -m '...archive_on_close: true'",
                      "",
                      "$ squad notes promote --ref squad-flight --dry-run",
                      "  Ref           Promoted  Archived  Skipped",
                      "  squad-flight         1         1        0",
                      "",
                      "$ squad notes promote --ref squad-flight",
                      "  squad-flight         1         1        0",
                      "  \u2713 Promotion complete.",
                      "",
                      "$ git ls-tree --name-only -r squad-state",
                      "  archive/squad-flight/<commit>.json",
                      "  promoted/squad-flight/<commit>.json"],
                     accent=GREEN, label="throwaway repo demo (verbatim)")
        add_card(slide, MARGIN_X + 6.8, CONTENT_TOP, 5.3, 4.2, "\u26A0",
                 "Honest caveat (found while running this demo)",
                 ["Backend writes a bulk note at refs/notes/squad",
                  "That blocks nested refs/notes/squad/{agent}",
                  "Demo used refs/notes/squad-flight to proceed",
                  "Documented per-agent namespace currently collides",
                  "",
                  "Real systems have rough edges \u2014 architects deserve them, not the demo polish."],
                 accent=YELLOW, head_size=14, body_size=12)
    add_content_slide(
        eyebrow="Promotion in motion",
        title="Flagged notes promote and archive into orphan state",
        subtitle="`promote_to_permanent` survives the branch; `archive_on_close` is kept for posterity.",
        source="run on a throwaway repo \u2014 the demo repo at files/two-layer-demo-repo reproduces this exactly",
        body=body,
        accent=GREEN,
    )


def s23_external_state():
    def body(slide):
        # TOP \u2014 the user pain that drove the feature (3 quick cards)
        pains = [
            ("\u26A0", "\u201CEvery branch switch is noisy\u201D",
             "git status keeps showing .squad/ files as modified/deleted because each branch has different team state."),
            ("\u2717", "\u201Cgit clean -fdx wiped our team\u201D",
             "On disk under .squad/ \u2192 nukable by any aggressive clean. Months of decisions and history, gone."),
            ("\u2298", "\u201CWe don\u2019t want .squad/ in PRs\u201D",
             "Even with .gitignore, reviewers ask \u2018what is this folder?\u2019 \u2014 the marker file is the only thing left in the repo."),
        ]
        n = len(pains)
        gap = 0.18
        w = (CONTENT_W - gap * (n - 1)) / n
        h = 1.85
        for i, (icon, head, desc) in enumerate(pains):
            x = MARGIN_X + i * (w + gap)
            add_rect(slide, x, CONTENT_TOP, w, h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, x, CONTENT_TOP, w, 0.08, fill=YELLOW, line=YELLOW)
            add_text_box(slide, x + 0.20, CONTENT_TOP + 0.18, 0.60, 0.45, icon,
                         size=22, color=YELLOW, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, x + 0.85, CONTENT_TOP + 0.18, w - 1.0, 0.45, head,
                         size=14, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, x + 0.20, CONTENT_TOP + 0.75, w - 0.40, h - 0.85, desc,
                         size=12, color=MUTED, font="Segoe UI")
        # MIDDLE \u2014 the answer (operator commands + use-when)
        mid_y = CONTENT_TOP + h + 0.20
        add_terminal(slide, MARGIN_X, mid_y, 6.5, 2.0,
                     ["$ squad externalize",
                      "  \u2713 moved 24 state files to:",
                      "    %APPDATA%\\squad\\projects\\my-monorepo\\",
                      "  \u2713 left .squad/config.json marker in repo",
                      "",
                      "$ squad internalize     # move state back into the repo"],
                     accent=ACCENT, label="the operator commands")
        add_card(slide, MARGIN_X + 6.7, mid_y, 5.4, 2.0, "\u2192",
                 "What you get",
                 ["Clean git status across all branches",
                  "State survives git clean / fresh clones",
                  "Per-repo global directory (no cross-talk)",
                  "Marker file stays in repo for re-link"],
                 accent=ACCENT, head_size=14, body_size=12)
        # BOTTOM \u2014 the honest caveat
        cap_y = mid_y + 2.20
        add_rect(slide, MARGIN_X, cap_y, CONTENT_W, 0.55, fill=PANEL2, line=YELLOW, corner=True)
        add_text_box(slide, MARGIN_X + 0.30, cap_y + 0.10, CONTENT_W - 0.60, 0.40,
                     "Key distinction: external state is per repo, per machine. "
                     "It\u2019s NOT a substitute for two-layer when you need shared team state across machines.",
                     size=12, color=YELLOW, bold=True, align=PP_ALIGN.CENTER, font="Segoe UI Semibold")
    add_content_slide(
        eyebrow="Why people asked for this",
        title="External state \u2014 when `.squad/` shouldn\u2019t even live in the repo",
        subtitle="Three real user pains drove this. It\u2019s an escape hatch \u2014 not a replacement for two-layer.",
        source="docs/.../features/external-state.md:65-109, 176-207 \u2022 issue threads driving it",
        body=body,
        notes=("Real motivation: power users with multiple checkouts / branch-heavy workflows kept getting bitten by "
               ".squad/ files appearing in git status on every switch. Externalize moves the noise out of the repo entirely "
               "by relocating mutable state to an OS-specific app-data directory and leaving only a marker file in .squad/. "
               "Crucial framing: this is a single-user / single-machine convenience. If you want teammates to share state, "
               "you want two-layer (shared via squad-state branch). The two features solve different problems."),
    )


# ====================== PART 3: SPAWNING ======================

def s24_section_spawning():
    add_section_divider(
        eyebrow="Pattern 3",
        title="\u2782 Spawning the\nright shape of team",
        lede="The scaling unit is a purpose-built workspace with a mission brief \u2014 not more agents in one chat.",
        accent=GREEN,
    )


def s25_fanout():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.5, 3.0,
                     ["User: \"Build auth\"",
                      "",
                      "Coordinator spawns in parallel:",
                      "  \u25CB Backend  : endpoints",
                      "  \u25CB Frontend : login form",
                      "  \u25CB Tester   : auth scenarios",
                      "  \u25CB Security : threat model"],
                     accent=GREEN, label="fan-out is the default")
        add_card(slide, MARGIN_X + 6.7, CONTENT_TOP, 5.4, 3.0, "\u2638",
                 "Serialize only for these",
                 ["Lead architecture decision gates implementation",
                  "Reviewer approval gates merge",
                  "Agent B needs Agent A's output file",
                  "User clarification is required"],
                 accent=GREEN, head_size=15, body_size=12)
    add_content_slide(
        eyebrow="Coordinator logic",
        title="Fan-out first \u2014 serialize only for real gates",
        subtitle="Parallel work is the default. Sync is for data dependencies and reviewer gates.",
        source="docs/.../features/parallel-execution.md:25-56, 86-99",
        body=body,
        accent=GREEN,
    )


def s26_worktrees():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.5, 3.0,
                     ["# issue #42 becomes an isolated workspace",
                      "git worktree add \\",
                      "  C:\\src\\squad-42 \\",
                      "  -b squad/42-fix-login",
                      "",
                      "# coordinator passes WORKTREE_PATH in the spawn prompt",
                      "# child does not switch branches inside the worktree"],
                     accent=GREEN, label="worktree lifecycle")
        add_card(slide, MARGIN_X + 6.7, CONTENT_TOP, 5.4, 3.0, "\u2691",
                 "What the child gets",
                 ["Dedicated branch and checkout",
                  "Known WORKTREE_PATH in the spawn prompt",
                  "No branch switching inside child",
                  "Build / test / commit happen in the worktree"],
                 accent=GREEN, head_size=15, body_size=12)
    add_content_slide(
        eyebrow="Isolation",
        title="Worktrees are the child workspace primitive",
        subtitle="A spawned mission needs a safe place to work without disrupting HQ.",
        source=".squad-templates/worktree-reference.md:48-127",
        body=body,
        accent=GREEN,
    )


def s27_subsquads():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.5, 3.0,
                     ["// .squad/workstreams.json",
                      "{",
                      "  \"workstreams\": [",
                      "    {",
                      "      \"name\": \"ui-team\",",
                      "      \"labelFilter\": \"team:ui\",",
                      "      \"folderScope\": [\"packages/ui\"]",
                      "    }",
                      "  ]",
                      "}"], accent=GREEN, label="stream config")
        add_card(slide, MARGIN_X + 6.7, CONTENT_TOP, 5.4, 3.0, "\u29C9",
                 "Why it scales",
                 ["Each Codespace sees only its label lane",
                  "Folder scope reduces accidental overlap",
                  "Branch-per-issue stays the real safety boundary",
                  "One machine can switch streams when cost matters"],
                 accent=GREEN, head_size=15, body_size=12)
    add_content_slide(
        eyebrow="Horizontal partitioning",
        title="SubSquads: one repo, many scoped teams",
        subtitle="Labels + folder scopes let multiple Squad instances pick up the right work.",
        source="packages/.../streams/resolver.ts:19-27, 124-181 \u2022 commands/streams.ts:1-48",
        body=body,
        accent=GREEN,
    )


def s28_hq_pattern():
    def body(slide):
        items = [
            ("\u2691", "HQ Squad",
             ["Owns the goal",
              "Chooses child shape",
              "Defines fan-in contract"], ACCENT),
            ("\u270E", "Mission brief",
             ["Scope, files, constraints",
              "Acceptance criteria",
              "Time / budget bounds"], YELLOW),
            ("\u29C9", "Child workspace",
             ["Worktree",
              "Tailored team shape",
              "Local mutable state"], GREEN),
            ("\u2705", "Structured result",
             ["PR / branch",
              "Decisions made",
              "Open risks"], VIOLET),
        ]
        cards_row(slide, CONTENT_TOP, 3.6, items)
        add_text_box(slide, MARGIN_X, CONTENT_TOP + 3.8, CONTENT_W, 0.7,
                     "Think \u2018mini project inside the project\u2019: a clear brief, isolated workspace, limited authority.",
                     size=15, color=MUTED, align=PP_ALIGN.CENTER)
    add_content_slide(
        eyebrow="Architecture pattern",
        title="HQ-and-child pattern",
        subtitle="A coordinator launches a child squad with a mission and collects a structured result.",
        source="Tamir Part 12 (fanout-squads) \u2022 .squad-templates/worktree-reference.md",
        body=body,
        accent=GREEN,
    )


def s29_mission_brief():
    def body(slide):
        # 2-column layout to prevent clipping of the Authority section
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 5.95, 4.4,
                     ["# Mission: Auth cache hardening",
                      "",
                      "Goal: reduce auth profile API latency",
                      "      without changing auth semantics.",
                      "",
                      "Scope:",
                      "  - packages/auth/**",
                      "  - tests/auth/**",
                      "",
                      "Out of scope:",
                      "  - packages/billing/**",
                      "  - infra/**"],
                     accent=GREEN, label="mission.md  \u2014  goal + scope")
        add_terminal(slide, MARGIN_X + 6.20, CONTENT_TOP, 5.95, 4.4,
                     ["Acceptance:",
                      "  - P95 < 80ms on /api/users/{id}",
                      "  - No new dependencies",
                      "  - All existing tests pass",
                      "",
                      "Authority:",
                      "  - May write under packages/auth,",
                      "                  tests/auth",
                      "  - Must NOT write secrets",
                      "  - Must NOT push to main",
                      "",
                      "Reporting:",
                      "  - Hand back PR URL + decisions made"],
                     accent=YELLOW, label="mission.md  \u2014  acceptance + authority")
    add_content_slide(
        eyebrow="Copy/paste artifact",
        title="Mission brief template \u2014 a contract, not a vague request",
        subtitle="Goal + scope on the left, acceptance + authority on the right. Two terminals so the Authority block never clips.",
        body=body,
        accent=GREEN,
    )


def s30_fanin_guardrails():
    def body(slide):
        items = [
            ("\u2638", "Authorized files",
             ["Tell the child what it may inspect/change"], ACCENT),
            ("\u29C9", "Branch discipline",
             ["No direct push to protected branches"], VIOLET),
            ("\u2698", "State protocol",
             ["Use runtime state tools",
              "Never git-write backend state by hand"], GREEN),
            ("\u270E", "Reviewer gates",
             ["Rejected author is locked out",
              "A different agent revises"], YELLOW),
            ("\u2691", "RAI pass",
             ["Rai catches secrets",
              "Catches unsafe patterns"], RED),
            ("\u2705", "Fan-in schema",
             ["status / branch / pr",
              "decisions, open risks"], GREEN),
        ]
        w = (CONTENT_W - 0.36) / 3
        h = 1.8
        for i, item in enumerate(items):
            icon, heading, bullets, color = item
            row, col = divmod(i, 3)
            x = MARGIN_X + col * (w + 0.18)
            y = CONTENT_TOP + row * (h + 0.18)
            add_card(slide, x, y, w, h, icon, heading, bullets,
                     accent=color, head_size=14, body_size=12)
    add_content_slide(
        eyebrow="Don't let parallelism become chaos",
        title="Guardrails for spawned squads",
        subtitle="The moment you spawn child teams, you need boundaries.",
        source="Reviewer Rejection Protocol \u2022 Rai charter \u2022 squad_state tools",
        body=body,
        accent=GREEN,
    )


# ====================== PART 4: CROSS-SQUAD ======================

def s31_section_cross():
    add_section_divider(
        eyebrow="Pattern 4",
        title="\u2783 Cross-squad\ncommunication",
        lede="A squad's public API is a manifest, a work request, and a clean feedback channel \u2014 not its internal history.",
        accent=YELLOW,
    )


def s32_manifest():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.6, 3.3,
                     ["// .squad/manifest.json",
                      "{",
                      "  \"name\": \"platform-squad\",",
                      "  \"version\": \"1.0.0\",",
                      "  \"description\": \"Platform engineering squad\",",
                      "  \"capabilities\": [\"infra\", \"observability\"],",
                      "  \"contact\": {",
                      "    \"repo\": \"acme/platform\",",
                      "    \"labels\": [\"squad:platform\"]",
                      "  },",
                      "  \"accepts\": [\"issues\"],",
                      "  \"skills\": [\"...\"]",
                      "}"], accent=YELLOW, label="public contract")
        add_card(slide, MARGIN_X + 6.8, CONTENT_TOP, 5.3, 3.3, "\u270E",
                 "Why manifests matter",
                 ["Consumers do not inspect internal .squad/ files",
                  "Capabilities are explicit",
                  "Contact + labels route work cleanly",
                  "`accepts` declares boundaries"],
                 accent=YELLOW, head_size=15, body_size=12)
    add_content_slide(
        eyebrow="Discovery",
        title="Manifest = public API contract",
        subtitle="Each squad publishes what it can do and how to reach it.",
        source="docs/features/cross-squad-orchestration.md:89-132 \u2022 cross-squad.ts:31-49",
        body=body,
        accent=YELLOW,
    )


def s33_discover_delegate():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.0, 3.5,
                     ["$ squad discover",
                      "Known Squads:",
                      "  backend-squad",
                      "    Capabilities: api, database",
                      "    Accepts:      issues",
                      "    Contact:      acme/backend  squad:backend",
                      "",
                      "$ squad delegate backend-squad \\",
                      "  \"Implement caching for /users API\"",
                      "  \u2192 creates issue in target repo",
                      "  \u2192 labels: squad:backend + squad:cross-squad"],
                     accent=YELLOW, label="operator view")
        add_card(slide, MARGIN_X + 6.2, CONTENT_TOP, 5.9, 3.5, "\u2716",
                 "Anti-patterns",
                 ["Never write directly into another squad's .squad/",
                  "Never depend on another squad's internals",
                  "Do not share secrets in handoffs",
                  "Always include acceptance criteria",
                  "Avoid circular delegation"],
                 accent=RED, head_size=15, body_size=12)
    add_content_slide(
        eyebrow="Handoff protocol",
        title="Discover \u2192 Delegate \u2192 Track",
        subtitle="Delegation is issue-based. The target squad receives context, labels, and acceptance criteria.",
        source="docs/features/cross-squad-orchestration.md:31-80 \u2022 packages/squad-sdk/src/runtime/cross-squad.ts:152-282",
        body=body,
        accent=YELLOW,
    )


def s34_pattern_compare():
    def body(slide):
        items = [
            ("\u2191", "Upstream",
             ["Parent \u2192 child",
              "Flows: policies, skills, wisdom",
              "Config: upstream.json"], ACCENT),
            ("\u21C4", "Cross-squad",
             ["Peer \u2192 owning squad",
              "Flows: work requests",
              "Config: manifest + issues"], YELLOW),
            ("\u29C9", "Distributed mesh",
             ["Peer \u21C4 peer",
              "Flows: state billboards",
              "Config: mesh.json + sync"], GREEN),
            ("\u29C0", "Export/import",
             ["Snapshot copy",
              "Flows: trained knowledge",
              "Config: squad export/import"], VIOLET),
        ]
        cards_row(slide, CONTENT_TOP, 3.7, items)
    add_content_slide(
        eyebrow="Pattern comparison",
        title="Choose the right communication pattern",
        subtitle="These mechanisms solve different coordination problems.",
        source="docs/.../distributed-mesh.md:213-257 \u2022 multiple-squads.md:17-27",
        body=body,
        accent=YELLOW,
    )


def s35_mesh():
    def body(slide):
        add_text_box(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 0.7,
                     "\u201CThe filesystem is the mesh. Git is how the mesh crosses machine boundaries.\u201D",
                     size=18, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        items = [
            ("1", "Zone 1 \u2014 Local",
             ["Same host / filesystem",
              "Direct file read"], ACCENT),
            ("2", "Zone 2 \u2014 Remote-trusted",
             ["Different host, same org",
              "git pull from shared repo"], GREEN),
            ("3", "Zone 3 \u2014 Remote-opaque",
             ["Different org / no shared auth",
              "curl published SUMMARY.md"], YELLOW),
        ]
        cards_row(slide, CONTENT_TOP + 0.9, 3.2, items)
    add_content_slide(
        eyebrow="Across machines",
        title="Distributed mesh: three trust zones",
        subtitle="Remote state must be materialized locally before agents read it.",
        source="docs/.../features/distributed-mesh.md:15-43 \u2022 SKILL.md:42-54",
        body=body,
        accent=YELLOW,
    )


# ====================== PART 5: MONOREPOS ======================

def s36a_section_monorepo():
    add_section_divider(
        eyebrow="Pattern 5",
        title="\u2784 Monorepo without\nstepping on yourself",
        lede="Multiple squads in one git repo: agent file at the root, .squad/ per package, workflows where they actually run.",
        accent=VIOLET,
        notes=("PR #939 + #1005 added monorepo subfolder support. "
               "The pain it solves: 'squad init' used to want git init, "
               "place workflows next to .squad/, and otherwise act like the "
               "subfolder was its own world. None of that works in a real monorepo."),
    )


def s36b_monorepo_topology():
    def body(slide):
        # Left: the layout
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.5, 4.0,
                     ["acme-monorepo/                       # git root",
                      "  .git/                              # the only .git",
                      "  .github/",
                      "    agents/squad.agent.md            # \u2190 single discovery point",
                      "    workflows/                       # \u2190 only here runs",
                      "      squad-ci.yml",
                      "      squad-heartbeat.yml",
                      "  packages/",
                      "    ui/",
                      "      .squad/                        # team-alpha",
                      "        team.md  routing.md  decisions.md",
                      "    api/",
                      "      .squad/                        # team-bravo",
                      "        team.md  routing.md  decisions.md",
                      "    infra/",
                      "      .squad/                        # team-charlie",
                      "        team.md  routing.md  decisions.md"],
                     accent=VIOLET, label="one repo, many squads, one agent file")
        add_card(slide, MARGIN_X + 6.7, CONTENT_TOP, 5.4, 4.0, "\u29C9",
                 "Two roots, one repo",
                 ["agentFileRoot  \u2192  git root",
                  "teamRoot       \u2192  package subfolder",
                  "Copilot discovers ONE agent file at the root",
                  "Each .squad/ owns its team identity, decisions, history",
                  "Workflows live at the git root \u2014 Actions ignores subfolders",
                  "Multiple squads do NOT collide on shared workflow files"],
                 accent=VIOLET, head_size=14, body_size=12)
    add_content_slide(
        eyebrow="Topology",
        title="One git root, many squads, one agent file",
        subtitle="Squad split its own concept of \u2018root\u2019 in two so monorepos can host independent teams.",
        source="packages/squad-sdk/src/config/init.ts:144-148, 1262, 1298-1305 \u2022 PR #939",
        body=body,
        accent=VIOLET,
    )


def s36c_monorepo_init():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 4.2,
                     ["$ cd acme-monorepo/packages/ui",
                      "$ squad init",
                      "",
                      "Detected parent git repo at acme-monorepo/  \u2192 monorepo subfolder mode",
                      "  \u2713 .squad/ created in packages/ui/",
                      "  \u2713 .github/agents/squad.agent.md placed at git root",
                      "  \u26A0 Skipped GitHub Actions workflows in monorepo-subfolder mode \u2014",
                      "    workflows must be at the git root. Set up workflows manually or",
                      "    use a single shared workflow for all squads.",
                      "",
                      "// behind the scenes (SDK contract)",
                      "initSquad({",
                      "  teamRoot:      \"packages/ui\",          // .squad/ lives here",
                      "  agentFileRoot: \"acme-monorepo\",        // git root \u2014 Copilot discovery",
                      "  includeWorkflows: true,                  // honored only at root",
                      "});",
                      "",
                      "// invariants (enforced by tests)",
                      "  packages/ui/.git/                         \u2192 does not exist (no nested init)",
                      "  acme-monorepo/.github/agents/squad...md   \u2192 exists",
                      "  packages/ui/.github/agents/squad...md     \u2192 does not exist"],
                     accent=VIOLET, label="init in a subfolder \u2014 verbatim shape")
    add_content_slide(
        eyebrow="Activation",
        title="`squad init` inside a package: what actually happens",
        subtitle="No nested git init, agent file at the root, workflows skipped with a clear warning.",
        source="packages/squad-cli/src/cli/core/init.ts:30-44 \u2022 test/init-scaffolding.test.ts:246-275",
        body=body,
        accent=VIOLET,
        notes=("Walk through the invariants. The test at init-scaffolding.test.ts:260-274 enforces every one: "
               "no nested .git, .squad/ in subfolder, agent file at root, no workflow files in subfolder, "
               "and a warnings entry explaining why."),
    )


# ====================== BONUS: PER-AGENT MODEL SELECTION ======================
# Came up from a real audience question: throttling on Opus 4.7 \u2192 needed to
# dial models per-agent. The 4-layer cascade IS architecture worth showing.

def s36d_model_cascade():
    def body(slide):
        # 4 stacked bands \u2014 same visual treatment as the 3-layer memory slide.
        layers = [
            {
                "tag":      "LAYER 0",
                "name":     "Persistent config",
                "subtitle": "Survives across sessions \u2022 user set it once \u2022 read on every spawn",
                "where":    ".squad/config.json",
                "example":  '"defaultModel": "claude-opus-4.8",  "agentModelOverrides": { "Scribe": "claude-haiku-4.5" }',
                "color":    VIOLET,
            },
            {
                "tag":      "LAYER 1",
                "name":     "Session directive",
                "subtitle": "Verbal \u2022 sticks until session ends or contradicted",
                "where":    "natural language at the prompt",
                "example":  "\u201Cuse claude-opus-4.8 for this session\u201D   \u201Csave costs\u201D",
                "color":    ACCENT,
            },
            {
                "tag":      "LAYER 2",
                "name":     "Charter preference",
                "subtitle": "Per-agent default in version control \u2022 lives with the agent identity",
                "where":    ".squad/agents/{name}/charter.md  \u2014 ## Model section",
                "example":  "Preferred: claude-opus-4.8",
                "color":    YELLOW,
            },
            {
                "tag":      "LAYER 3",
                "name":     "Task-aware auto-select",
                "subtitle": "Cost-first \u2022 unless code or prompt architecture is being written",
                "where":    "role + task-output tables",
                "example":  "code\u2192sonnet  \u2022  prompts\u2192sonnet  \u2022  docs/triage/logs\u2192haiku  \u2022  vision\u2192opus",
                "color":    GREEN,
            },
        ]
        band_h = 1.10
        gap    = 0.12
        for i, layer in enumerate(layers):
            y = CONTENT_TOP + i * (band_h + gap)
            color = layer["color"]
            add_rect(slide, MARGIN_X, y, CONTENT_W, band_h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, MARGIN_X, y, 0.18, band_h, fill=color, line=color)
            # Layer tag (top-left of band, vertical accent)
            tag_x = MARGIN_X + 0.35
            add_text_box(slide, tag_x, y + 0.10, 1.20, 0.30, layer["tag"],
                         size=11, color=color, bold=True, font="Consolas")
            # Name (next to tag, same row)
            add_text_box(slide, tag_x + 1.20, y + 0.05, 4.5, 0.42, layer["name"],
                         size=17, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            # Subtitle (below name)
            add_text_box(slide, tag_x, y + 0.50, 6.0, 0.30, layer["subtitle"],
                         size=10, color=MUTED, font="Segoe UI")
            # Where (third row, mono)
            add_text_box(slide, tag_x, y + 0.80, 6.0, 0.28, layer["where"],
                         size=10, color=DIM, font="Consolas")
            # Example column (right side)
            ex_x = MARGIN_X + 6.30
            ex_w = CONTENT_W - (ex_x - MARGIN_X) - 0.20
            add_text_box(slide, ex_x, y + 0.10, ex_w, 0.25, "EXAMPLE",
                         size=9, color=color, bold=True, font="Segoe UI Semibold")
            add_text_box(slide, ex_x, y + 0.35, ex_w, band_h - 0.40, layer["example"],
                         size=11, color=TEXT, font="Consolas")
    add_content_slide(
        eyebrow="Per-agent control",
        title="Model selection cascade \u2014 four layers, first match wins",
        subtitle="Check top-down, stop at the first hit. Each layer has a different scope: cross-session \u2192 session \u2192 charter \u2192 task.",
        source=".squad-templates/model-selection-reference.md:5-46 \u2022 .github/agents/squad.agent.md:414-420",
        body=body,
        accent=VIOLET,
        notes=("Real audience scenario: Opus 4.7 throttling \u2192 needs to fail over to Opus 4.8 cleanly. "
               "Persistent config is the lever \u2014 set defaultModel once, override Scribe/Ralph to haiku to preserve quota. "
               "Anchor: 'each layer has a different scope (cross-session \u2192 session \u2192 charter \u2192 task) \u2014 "
               "that's the architectural insight, not the cascade itself.'"),
    )


def s36e_model_recipe():
    def body(slide):
        # Left: the variant-string insight
        add_card(slide, MARGIN_X, CONTENT_TOP, 5.7, 4.3, "\u2698",
                 "The model NAME is the knob",
                 ["Effort and context window are encoded in the variant string",
                  "claude-opus-4.7         \u2192 standard",
                  "claude-opus-4.7-high    \u2192 high reasoning effort",
                  "claude-opus-4.7-xhigh   \u2192 extra-high effort",
                  "claude-opus-4.7-1m-internal \u2192 1M context window",
                  "No separate 'effort' or 'context' API \u2014 the SDK just receives the string",
                  "When new variants ship, swap the string \u2014 no agent rewrites"],
                 accent=ACCENT, head_size=14, body_size=12)
        # Right: practical recipe terminal
        add_terminal(slide, MARGIN_X + 5.9, CONTENT_TOP, 6.2, 4.3,
                     ["// .squad/config.json  \u2014 survive throttling cleanly",
                      "{",
                      "  \"version\": 1,",
                      "  \"defaultModel\": \"claude-opus-4.8\",",
                      "  \"agentModelOverrides\": {",
                      "    \"Scribe\":   \"claude-haiku-4.5\",   // logs, merges \u2014 cheap",
                      "    \"Ralph\":    \"claude-haiku-4.5\",   // backlog loop \u2014 cheap",
                      "    \"Lead\":     \"claude-opus-4.8\",    // architecture \u2014 premium",
                      "    \"Designer\": \"claude-opus-4.5\"     // vision required",
                      "  }",
                      "}",
                      "",
                      "// or per session, no file edit:",
                      "> \u201Cuse claude-opus-4.8 for this session\u201D",
                      "",
                      "// or pin in a charter (lives with the agent):",
                      "// .squad/agents/lead/charter.md",
                      "## Model",
                      "Preferred: claude-opus-4.8"],
                     accent=GREEN, label="practical recipe \u2014 cost first, code second")
    add_content_slide(
        eyebrow="Practical recipe",
        title="Surviving throttling: dial premium where it matters, haiku everywhere else",
        subtitle="One config file. Three escape hatches when you need to deviate.",
        source=".squad-templates/model-selection-reference.md:7-11, 26-46",
        body=body,
        accent=GREEN,
        notes=("Anchor: Scribe and Ralph do mechanical work \u2014 they should NEVER be on a premium model. "
               "Lead and Designer get bumped because their output gates merge / requires vision. "
               "The variant-string design means adding 4.8-high / 4.8-1m later is a single-line config edit."),
    )


# ====================== WORKFLOW: PROMPT \u2192 NOTE ======================
# A progressive-reveal sequence (Aspire-style). Same flow strip at the top of
# every slide, one stage active per slide, detail panel below changes.

WORKFLOW_NODES = [
    "1. User prompt",
    "2. Coordinator + agents",
    "3. Decision emerges",
    "4. Classify + call tool",
    "5. Write to state",
    "6. Verify",
]


def _workflow_header(slide, *, active, accent=ACCENT):
    add_flow_strip_highlighted(slide, MARGIN_X, CONTENT_TOP - 0.05, CONTENT_W, 0.85,
                               WORKFLOW_NODES, active=active, accent=accent)


def s36f_synthesis_divider():
    add_section_divider(
        eyebrow="Synthesis",
        title="Now watch all five patterns\nfire in sequence",
        lede="One natural-language directive \u2192 classifier \u2192 governed write \u2192 git note \u2192 orphan branch. Six stages, real artifacts.",
        accent=ACCENT,
        notes=("This is the climax of the talk. After 5 isolated patterns, the next 8 slides show a single sentence "
               "(\u201CAlways run the SDK tests before merging.\u201D) trace through every layer we built. "
               "Don\u2019t rush this section \u2014 it\u2019s the payoff."),
    )


def s37_workflow_overview():
    def body(slide):
        _workflow_header(slide, active=None)
        # Description card
        add_rect(slide, MARGIN_X, CONTENT_TOP + 1.05, CONTENT_W, 3.4, fill=PANEL2, line=BORDER, corner=True)
        add_text_box(slide, MARGIN_X + 0.3, CONTENT_TOP + 1.20, CONTENT_W - 0.6, 0.5,
                     "What this sequence shows",
                     size=18, color=ACCENT, bold=True)
        add_text_box(slide, MARGIN_X + 0.3, CONTENT_TOP + 1.75, CONTENT_W - 0.6, 2.6,
                     "We follow ONE durable sentence \u2014 \u201CAlways run the SDK tests before merging.\u201D \u2014 "
                     "from the moment the user types it through to the git note and orphan branch where it lives.\n\n"
                     "Each of the next five slides zooms one stage. Real code from packages/squad-sdk/src/memory "
                     "and packages/squad-sdk/src/state-backend.ts; real commands you can run on the deck repo afterwards.",
                     size=15, color=MUTED)
    add_content_slide(
        eyebrow="End-to-end sequence",
        title="One sentence, six stages, real git artifacts",
        subtitle="Trace a directive from natural language through to the orphan branch and git note.",
        source="memory/index.ts \u2022 state-backend.ts \u2022 .github/agents/squad.agent.md (Directive Capture)",
        body=body,
    )


def s38_workflow_step1_prompt():
    def body(slide):
        _workflow_header(slide, active=0)
        # Detail
        add_terminal(slide, MARGIN_X, CONTENT_TOP + 1.05, 6.5, 3.4,
                     ["> Always run the SDK tests before merging.",
                      "",
                      "(Plain natural-language directive.",
                      " No command. No flag. No JSON.",
                      " The user does not say \u201Cremember this\u201D.)"],
                     accent=ACCENT, label="what the human typed")
        add_card(slide, MARGIN_X + 6.7, CONTENT_TOP + 1.05, 5.4, 3.4, "\u270E",
                 "Why this is hard",
                 ["The team has no way to know what's important",
                  "Tomorrow another agent will violate this",
                  "Old solution: tell the user to edit decisions.md",
                  "New solution: the coordinator listens for directive signals"],
                 accent=ACCENT, head_size=14, body_size=12)
    add_content_slide(
        eyebrow="Stage 1 \u2014 User prompt",
        title="A sentence the team should remember forever",
        subtitle="Captured before any work is dispatched. This is the input to the whole pipeline.",
        source=".github/agents/squad.agent.md \u2014 Directive Capture",
        body=body,
    )


def s39_workflow_step2_coordinator():
    def body(slide):
        _workflow_header(slide, active=1)
        add_terminal(slide, MARGIN_X, CONTENT_TOP + 1.05, 6.5, 3.4,
                     ["Coordinator (Squad agent):",
                      "  pattern check: /^\\s*(always|never|must|do not)\\b/i",
                      "  hit \u2192 \u201Calways\u201D \u2192 treat as DIRECTIVE",
                      "",
                      "Routing branch: directive-capture",
                      "  \u2022 do NOT spawn implementation agents",
                      "  \u2022 do NOT route to Lead for review",
                      "  \u2022 DO call governed memory tool now",
                      "",
                      "(Scribe is informed; no fan-out happens.)"],
                     accent=ACCENT, label="coordinator detects + decides")
        add_card(slide, MARGIN_X + 6.7, CONTENT_TOP + 1.05, 5.4, 3.4, "\u2638",
                 "When agents WOULD chat",
                 ["Implementation requests fan out to specialists",
                  "Directive capture does NOT \u2014 it's a pre-routing hook",
                  "Pure knowledge events skip the team chatter",
                  "Same pipeline shape; different participants per request"],
                 accent=GREEN, head_size=14, body_size=12)
    add_content_slide(
        eyebrow="Stage 2 \u2014 Coordinator + agents",
        title="The coordinator chooses: knowledge event vs. work",
        subtitle="A directive is captured before any specialist is spawned.",
        source=".github/agents/squad.agent.md \u2014 Directive Capture \u2022 routing table",
        body=body,
    )


def s40_workflow_step3_decision():
    def body(slide):
        _workflow_header(slide, active=2)
        add_terminal(slide, MARGIN_X, CONTENT_TOP + 1.05, CONTENT_W, 3.4,
                     ["// candidate knowledge object (in coordinator's head)",
                      "{",
                      "  title:    \"Run SDK tests before merging\",",
                      "  content:  \"Always run the SDK tests before merging.\",",
                      "  author:   \"squad-coordinator\",",
                      "  metadata: { source: \"user-directive\" }",
                      "}",
                      "",
                      "// the coordinator has decided: this should be durable team memory.",
                      "// it has NOT yet decided WHERE it goes or HOW it loads.",
                      "// that's the next stage \u2014 a tool call to memory.classify()."],
                     accent=YELLOW, label="durable fact \u2014 not yet classified, not yet stored")
    add_content_slide(
        eyebrow="Stage 3 \u2014 Decision emerges",
        title="A candidate fact, not yet a memory",
        subtitle="The coordinator constructs the request object. The classifier still has to vet it.",
        source="memory/index.ts:667-680 \u2014 MemoryWriteRequest shape",
        body=body,
        accent=YELLOW,
    )


def s41_workflow_step4_classify():
    def body(slide):
        _workflow_header(slide, active=3, accent=GREEN)
        add_terminal(slide, MARGIN_X, CONTENT_TOP + 1.05, 7.0, 3.4,
                     ["// memory/index.ts:594-605 \u2014 FORBIDDEN scan first (incl. CI/PR/build)",
                      "for ({ pattern, reason } of FORBIDDEN_PATTERNS) {",
                      "  if (pattern.test(content)) return FORBIDDEN;",
                      "}",
                      "",
                      "// memory/index.ts:610-619 \u2014 heuristics (TRANSIENT branch is",
                      "// unreachable today \u2014 same regex sits in FORBIDDEN_PATTERNS)",
                      "if (/\\b(CI|PR|build) (status|failed|...)\\b/i.test(content))",
                      "  memoryClass = 'TRANSIENT';   // dead code via FORBIDDEN overlap",
                      "else if (/^\\s*(always|never|must|do not)\\b/i.test(content))",
                      "  memoryClass = 'POLICY';      //  \u2190 our case fires here",
                      "else if (/\\b(decision|decided|...)\\b/i.test(content))",
                      "  memoryClass = 'DECISION';",
                      "else memoryClass = 'LOCAL';     // fallthrough",
                      "",
                      "// memory/index.ts:640 \u2014 load guidance from class",
                      "loadGuidance = loadGuidanceFor('POLICY');  //  \u2192 ALWAYS"],
                     accent=GREEN, label="real classifier code \u2014 with the honest overlap noted")
        add_terminal(slide, MARGIN_X + 7.2, CONTENT_TOP + 1.05, 4.9, 3.4,
                     ["$ squad memory classify \\",
                      "    \"Always run the SDK tests \\",
                      "     before merging.\"",
                      "",
                      "{",
                      "  \"class\":         \"POLICY\",",
                      "  \"allowed\":       true,",
                      "  \"destination\":   \"decisions/inbox\",",
                      "  \"loadGuidance\":  \"ALWAYS\",",
                      "  \"reason\":        \"Content is allowed\",",
                      "  ...",
                      "}"],
                     accent=ACCENT, label="actual CLI output")
    add_content_slide(
        eyebrow="Stage 4 \u2014 Classify + call tool",
        title="The classifier picks POLICY \u2192 loads ALWAYS",
        subtitle="The class is computed from real regex rules. Result becomes the routing key for storage.",
        source="packages/squad-sdk/src/memory/index.ts:588-665 \u2014 classify() heuristics + FORBIDDEN scan",
        body=body,
        accent=GREEN,
    )


def s42_workflow_step5_write():
    def body(slide):
        _workflow_header(slide, active=4, accent=VIOLET)
        # Left: SDK call chain
        add_terminal(slide, MARGIN_X, CONTENT_TOP + 1.05, 6.4, 3.4,
                     ["// SDK call chain (verbatim from memory/index.ts)",
                      "store.write(request)",
                      "  \u2192 classification = await this.classify(request)",
                      "  \u2192 id = randomUUID()",
                      "  \u2192 relPath = destinationPath('POLICY', id, title, author)",
                      "  \u2192 content = renderMemoryFile(id, 'POLICY', title, request)",
                      "  \u2192 storage.write(fullPath, content)            //  \u2190 (A)",
                      "  \u2192 index.push(entry); writeIndex(index)        //  \u2190 (B)",
                      "  \u2192 audit({ action: 'write', class: 'POLICY' }) //  \u2190 (C)",
                      "",
                      "// (A) the storage interface is the seam:",
                      "//     storage = new StateBackendStorageAdapter(twoLayerBackend, dir)",
                      "//     state-backend.ts:798-800 \u2014 routes write \u2192 backend.write()"],
                     accent=VIOLET, label="memory store \u2192 storage adapter")
        # Right: two-layer fan-out
        add_terminal(slide, MARGIN_X + 6.6, CONTENT_TOP + 1.05, 5.5, 3.4,
                     ["// TwoLayerBackend.write(relPath, content)",
                      "//   \u2190 state-backend.ts (composed)",
                      "",
                      "// (1) orphan branch  \u2014 durable",
                      "git update-index --add  squad-state",
                      "git commit-tree         squad-state",
                      "git update-ref refs/heads/squad-state",
                      "",
                      "// (2) git note       \u2014 best-effort, commit-scoped",
                      "git notes --ref=refs/notes/squad add \\",
                      "  -F <classification-json>",
                      "",
                      "// (3) audit log",
                      ".squad/memory/audit.jsonl  \u2190 redacted entry"],
                     accent=GREEN, label="two layers, one logical write")
    add_content_slide(
        eyebrow="Stage 5 \u2014 Write to state",
        title="The tool stores it in the right place \u2014 actually",
        subtitle="One logical write fans out: orphan branch (durable), git note (commit-scoped), audit (forensics).",
        source="memory/index.ts:667-823 \u2022 state-backend.ts:791-840 \u2022 TwoLayerBackend (composed)",
        body=body,
        accent=VIOLET,
    )


def s43_workflow_step6_verify():
    def body(slide):
        _workflow_header(slide, active=5, accent=GREEN)
        add_terminal(slide, MARGIN_X, CONTENT_TOP + 1.05, CONTENT_W, 3.7,
                     ["# (1) working tree is still clean \u2014 no decisions diff in your PR",
                      "$ git status -s",
                      "",
                      "# (2) durable store on the orphan branch carries the fact",
                      "$ git show squad-state:decisions/inbox/coordinator-run-sdk-tests-*.md",
                      "  ---",
                      "  class: POLICY",
                      "  loadGuidance: ALWAYS",
                      "  ---",
                      "  Always run the SDK tests before merging.",
                      "",
                      "# (3) the commit you JUST made carries the classification as a note",
                      "$ git notes --ref=refs/notes/squad show HEAD",
                      "  { \"class\": \"POLICY\", \"loadGuidance\": \"ALWAYS\", \"id\": \"e96f...\" }",
                      "",
                      "# (4) audit log records who wrote it and why",
                      "$ squad memory audit | tail -1",
                      "  { \"action\":\"write\", \"class\":\"POLICY\", \"actor\":\"coordinator\", ... }"],
                     accent=GREEN, label="four commands. four real artifacts.")
    add_content_slide(
        eyebrow="Stage 6 \u2014 Verify",
        title="Four commands prove the write actually landed",
        subtitle="Working tree, orphan branch, git note, audit log \u2014 each tells a different part of the truth.",
        source="git status \u2022 git show squad-state \u2022 git notes refs/notes/squad \u2022 squad memory audit",
        body=body,
        accent=GREEN,
        notes=("Run these on the deck repo. Each one returns real content because two-layer is enabled. "
               "Branch links to share live: github.com/.../tree/squad-state and "
               "github.com/.../tree/squad/advanced-squad-session-slides."),
    )


def s44_workflow_recap():
    def body(slide):
        # Show full strip un-dimmed
        add_flow_strip(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 0.85, WORKFLOW_NODES, accent=ACCENT)
        # Recap row
        items = [
            ("\u270E", "Prompt",         "Plain natural language. The trigger.",        ACCENT),
            ("\u2638", "Coordinator",    "Directive signal \u2192 pre-routing hook.",   ACCENT),
            ("\u2691", "Decision",       "Candidate fact \u2014 not yet a memory.",     YELLOW),
            ("\u26A1", "Classify",       "Real regex \u2192 POLICY \u2192 ALWAYS.",      GREEN),
            ("\u2698", "Write",          "Orphan branch + note + audit.",                VIOLET),
            ("\u2705", "Verify",         "Four commands. Four real artifacts.",          GREEN),
        ]
        n = len(items)
        w = (CONTENT_W - 0.18 * (n - 1)) / n
        h = 2.9
        for i, item in enumerate(items):
            icon, heading, body_line, color = item
            x = MARGIN_X + i * (w + 0.18)
            y = CONTENT_TOP + 1.05
            add_rect(slide, x, y, w, h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, x, y, w, 0.08, fill=color, line=color)
            add_text_box(slide, x + 0.1, y + 0.18, w - 0.2, 0.5,
                         icon, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
            add_text_box(slide, x + 0.1, y + 0.75, w - 0.2, 0.5,
                         heading, size=15, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
            add_text_box(slide, x + 0.15, y + 1.35, w - 0.3, 1.4,
                         body_line, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_content_slide(
        eyebrow="Recap",
        title="Six stages, one durable sentence",
        subtitle="Memory becomes infrastructure when each stage is observable and replayable.",
        body=body,
    )


def s37_maturity():
    def body(slide):
        levels = [
            ("L1", "One Squad",            "Use memory + decisions"),
            ("L2", "Reusable team shapes", "Add presets"),
            ("L3", "Parallel missions",    "Worktrees + SubSquads"),
            ("L4", "Many squads",          "Manifest discovery + delegation"),
            ("L5", "Distributed org",      "Mesh + upstream + policy cascade"),
        ]
        n = len(levels)
        w = (CONTENT_W - 0.18 * (n - 1)) / n
        h = 3.5
        for i, (lvl, head, body_line) in enumerate(levels):
            x = MARGIN_X + i * (w + 0.18)
            add_rect(slide, x, CONTENT_TOP, w, h, fill=PANEL, line=BORDER, corner=True)
            add_rect(slide, x, CONTENT_TOP, w, 0.08, fill=ACCENT, line=ACCENT)
            add_text_box(slide, x + 0.1, CONTENT_TOP + 0.25, w - 0.2, 0.5,
                         lvl, size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
            add_text_box(slide, x + 0.1, CONTENT_TOP + 0.85, w - 0.2, 0.6,
                         head, size=15, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
            add_text_box(slide, x + 0.15, CONTENT_TOP + 1.55, w - 0.3, 1.8,
                         body_line, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_content_slide(
        eyebrow="Maturity ladder",
        title="Advanced operating model",
        subtitle="Don't turn everything on at once. Add patterns when the pain appears.",
        body=body,
    )


def s37b_human_payoff():
    """Closing synthesis slide — answers slide 02's three questions in human terms,
    not in git-log terms. Per tamresearch1 squad review (Picard/Troi/Guinan): the deck
    showed the mechanism honestly; this slide shows the RELIEF."""
    def body(slide):
        # Three answers, each in plain human language
        items = [
            ("\u2698", "How does it remember?",
             ["Agents stop asking what they",
              "already answered.",
              "",
              "\u201CWe decided to use Postgres\u201D \u2014",
              "the team\u2019s shared brain reloads it",
              "into every future spawn,",
              "no matter who wrote it."],
             ACCENT),
            ("\u29C9", "How do we change team shape?",
             ["A mission becomes a workspace,",
              "not a chat thread.",
              "",
              "Spawn a child squad with a brief,",
              "let it work in its own worktree,",
              "collect a clean PR back \u2014",
              "without disrupting HQ."],
             GREEN),
            ("\u21C4", "How do squads talk?",
             ["A manifest, an issue, a verdict.",
              "",
              "No agent of yours ever pokes",
              "another team\u2019s .squad/ directly.",
              "Public contract in, structured",
              "result out. Just like a service."],
             YELLOW),
        ]
        cards_row(slide, CONTENT_TOP, 3.7, items)
        # Bottom: the relief sentence
        cap_y = CONTENT_TOP + 3.95
        add_rect(slide, MARGIN_X, cap_y, CONTENT_W, 0.85, fill=PANEL2, line=VIOLET, corner=True)
        add_text_box(slide, MARGIN_X + 0.30, cap_y + 0.10, CONTENT_W - 0.60, 0.30,
                     "What changes six months in",
                     size=12, color=VIOLET, bold=True, font="Segoe UI Semibold")
        add_text_box(slide, MARGIN_X + 0.30, cap_y + 0.42, CONTENT_W - 0.60, 0.40,
                     "git log stays clean. PRs are reviewable. Decisions you made on day 12 are still alive on day 180. "
                     "Your team doesn\u2019t lose context to keep the repo clean \u2014 it gets both.",
                     size=14, color=TEXT, align=PP_ALIGN.CENTER)
    add_content_slide(
        eyebrow="The payoff",
        title="The three questions, answered in human terms",
        subtitle="The talk opened with three questions. Here\u2019s what changes for your team when all five patterns are running.",
        body=body,
        notes=("This is the landing the talk was missing. Slide 02 asked three questions; this slide answers them in plain "
               "human language, not git-log terms. Lead in by reminding the audience of the opening questions, then walk "
               "the three cards. Land on the bottom callout: 'your team doesn\u2019t lose context to keep the repo clean \u2014 "
               "it gets both.' This is the line that makes the architecture feel like relief, not just mechanism."),
    )


def s38_cheatsheet():
    def body(slide):
        add_terminal(slide, MARGIN_X, CONTENT_TOP, 6.0, 4.0,
                     ["# memory",
                      "squad memory classify \"Always...\"",
                      "squad memory write --content \"...\" --class DECISION",
                      "squad memory search --query \"...\"",
                      "squad memory audit",
                      "squad memory provider             # check what's wired",
                      "",
                      "# state backends",
                      "squad init --state-backend two-layer",
                      "squad upgrade --state-backend two-layer",
                      "squad notes promote --ref squad/scribe",
                      "",
                      "# verify your setup",
                      "squad doctor                       # use FIRST when stuck"],
                     accent=ACCENT, label="local team")
        add_terminal(slide, MARGIN_X + 6.2, CONTENT_TOP, 6.0, 4.0,
                     ["# presets",
                      "squad preset init [--remote]",
                      "squad preset list / show / apply / save",
                      "",
                      "# spawning / scale",
                      "export SQUAD_TEAM=ui-team          # bash / zsh",
                      "$env:SQUAD_TEAM = \"ui-team\"        # PowerShell",
                      "squad subsquads list / activate ui-team",
                      "git worktree add ../squad-42 -b squad/42-...",
                      "",
                      "# cross-squad",
                      "squad discover",
                      "squad delegate <target> \"...\"",
                      "",
                      "# monorepo (run from inside a package)",
                      "cd packages/ui && squad init   # auto-detects subfolder mode"],
                     accent=GREEN, label="multi-team")
    add_content_slide(
        eyebrow="Leave-behind",
        title="Command cheat sheet",
        subtitle="The commands you're most likely to use after this session.",
        body=body,
    )


def s39_close():
    def body(slide):
        items = [
            ("\u2698", "Memory",
             ["Put facts in the right layer",
              "Class on write \u2192 guidance on read"], ACCENT),
            ("\u25A4", "State",
             ["Two-layer keeps PRs clean",
              "Orphan branch + notes do the heavy lifting"], VIOLET),
            ("\u2691", "Spawning",
             ["Give missions their own workspace",
              "Worktrees + mission briefs"], YELLOW),
            ("\u21C4", "Cross-squad",
             ["Use public contracts, not internals",
              "Manifest + issue handoff"], GREEN),
            ("\u29C9", "Monorepo",
             ["agentFileRoot vs teamRoot",
              "One agent file, many squads"], VIOLET),
        ]
        n = len(items)
        gap = 0.18
        w = (CONTENT_W - gap * (n - 1)) / n
        h = 3.4
        for i, (icon, heading, bullets, color) in enumerate(items):
            x = MARGIN_X + i * (w + gap)
            add_card(slide, x, CONTENT_TOP, w, h, icon, heading, bullets,
                     accent=color, head_size=14, body_size=12)
        add_text_box(slide, MARGIN_X, CONTENT_TOP + 3.7, CONTENT_W, 0.6,
                     "Thank you \u2014 questions?",
                     size=26, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, MARGIN_X, CONTENT_TOP + 4.25, CONTENT_W, 0.5,
                     "github.com/bradygaster/squad  \u2022  tamirdresher.com",
                     size=15, color=ACCENT, align=PP_ALIGN.CENTER)
    add_content_slide(
        eyebrow="Close",
        title="The mental model to remember",
        subtitle="Advanced Squad is operational architecture for AI teams \u2014 not 'more agents.'",
        body=body,
    )


def s40_appendix():
    def body(slide):
        lines = [
            "\u2022 package.json:3  \u2022  CHANGELOG.md:7-66  \u2014  Squad 0.10.0 scope",
            "\u2022 docs/.../features/memory.md:28-74, 113-166  \u2014  memory layers, governed memory commands",
            "\u2022 docs/.../concepts/memory-and-knowledge.md:28-79  \u2014  three-layer memory model",
            "\u2022 packages/squad-sdk/src/memory/index.ts:5-13, 428-440  \u2014  MemoryClass + loadGuidanceFor",
            "\u2022 packages/squad-sdk/src/memory/index.ts:863-950  \u2014  search() with class-aware filtering",
            "\u2022 docs/.../features/state-backends.md:21-156  \u2014  local/orphan/two-layer comparison",
            "\u2022 packages/squad-sdk/src/state-backend.ts:910-1186  \u2014  TwoLayerBackend + resolveStateBackend",
            "\u2022 packages/squad-cli/src/cli/commands/notes.ts  \u2014  squad notes promote",
            "\u2022 docs/.../features/external-state.md:65-188  \u2014  external state command path",
            "\u2022 packages/.../preset.ts:14-389  \u2014  preset init/list/show/apply/save (incl --remote)",
            "\u2022 packages/.../streams/resolver.ts:124-181  \u2014  SubSquads / workstream partitioning",
            "\u2022 docs/features/cross-squad-orchestration.md:31-132  \u2014  discover / delegate / manifest",
            "\u2022 packages/squad-sdk/src/runtime/cross-squad.ts:31-282  \u2014  manifest schema, discovery, delegation",
            "\u2022 docs/.../features/distributed-mesh.md:15-257  \u2014  mesh zones, lifecycle, examples",
            "\u2022 packages/squad-sdk/src/config/init.ts:144-148, 1262, 1298-1305  \u2014  monorepo subfolder support",
            "\u2022 packages/squad-cli/src/cli/core/init.ts:30-44  \u2014  detectParentGitRepo (monorepo)",
            "\u2022 test/init-scaffolding.test.ts:246-275  \u2014  monorepo invariants (no nested .git, agent at root)",
            "\u2022 PR #1004, #1035, #1145, #1022, #939, #1005; issues #1036, #1037  \u2014  the work behind this session",
            "\u2022 tamirdresher.com Part 7, 7b, 11, 12, 13  \u2014  the blog series this deck draws from",
        ]
        tb = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(CONTENT_TOP), Inches(CONTENT_W), Inches(4.5))
        tf = tb.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(2)
            run = p.add_run()
            run.text = line
            run.font.name = "Consolas"; run.font.size = Pt(11); run.font.color.rgb = MUTED
    add_content_slide(
        eyebrow="Grounding",
        title="Source appendix",
        subtitle="Every product claim above is grounded in accessible repo files.",
        body=body,
    )


# =================================================================
# BUILD
# =================================================================
s01_title()
s02_promise()
s03_roadmap()
s04_release()
s05_vocab()
s06_scaling_shift()

s07_section_memory()
s08_three_layers()
s09_governed_classes()
s10a_write_architecture()
s10b_memory_providers()
s09b_memory_tools()
s09c_classify_spotlight()
s10_retrieval()
s11_memory_tools_proof()
s12_session_evidence()
s13_directive_capture()
s14_memory_antipatterns()

s15_section_state()
s16_state_vs_code()
s17_backends_compare()
s18_two_layer_thesis()
s19_activation()
s20_live_upgrade()
s21_live_state_proof()
s21c_session_evidence()
s21c2_branch_topology()
s21c3_timeline()
s21c4_landings()
s21b_steady_state_safety_net()
s21d_pre_commit_fires()
s21e_safety_net_recovery()
s22_two_layer_demo_promotion()
s23_external_state()

s24_section_spawning()
s25_fanout()
s26_worktrees()
s27_subsquads()
s28_hq_pattern()
s29_mission_brief()
s30_fanin_guardrails()

s31_section_cross()
s32_manifest()
s33_discover_delegate()
s34_pattern_compare()
s35_mesh()

# Pattern 5 \u2014 monorepo subfolder mode (NEW)
s36a_section_monorepo()
s36b_monorepo_topology()
s36c_monorepo_init()

# Bonus \u2014 per-agent model selection (NEW)
s36d_model_cascade()
s36e_model_recipe()

# Synthesis section divider before the end-to-end workflow
s36f_synthesis_divider()

# End-to-end workflow \u2014 progressive reveal across 7 slides (NEW)
s37_workflow_overview()
s38_workflow_step1_prompt()
s39_workflow_step2_coordinator()
s40_workflow_step3_decision()
s41_workflow_step4_classify()
s42_workflow_step5_write()
s43_workflow_step6_verify()
s44_workflow_recap()

s37_maturity()
s37b_human_payoff()
s38_cheatsheet()
s40_appendix()
s39_close()

out = Path(r"C:\Users\tamirdresher\source\repos\squad-advanced-squad-session-slides\docs\advanced-squad-session\advanced-squad-session.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"slides={len(prs.slides)}")
print(out)
